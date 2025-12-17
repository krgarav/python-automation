# inspiration_slides.py

import os
import copy
from pptx import Presentation
from pptx.util import Emu
from PIL import Image, ImageOps
from pptx.enum.shapes import MSO_SHAPE_TYPE


# -----------------------------------------------------
# SPACE → Explanation field mapping
# -----------------------------------------------------
SPACE_EXPLANATION_MAP = {
    "Other": {"explanation": "short_textzp4lo2ts"},
    "Bathroom": {"explanation": "short_textyykqgp9s"},
    "Bedroom": {"explanation": "short_text79mzvngh"},
    "Dining Room": {"explanation": "short_textp906n78y"},
    "Entrance/Foyer/Hallway": {"explanation": "short_textsyh8d82l"},
    "Kitchen": {"explanation": "short_text3aufbldq"},
    "Living Room": {"explanation": "short_textrhhwmi8e"},
    "Outdoor area": {"explanation": "short_textbizb9r6q"},
    "Storage space": {"explanation": "short_textmgvse9ua"},
    "Study /Library room": {"explanation": "short_text6rb6983p"},
    "Toilet / powder room": {"explanation": "short_textnx6moc35"},
    "Entertainment": {"explanation": "short_textbm7l60t4"},
    "Lobby / Reception": {"explanation": "short_text0hb93h0z"},
    "Bar": {"explanation": "short_texthtjb2m8p"},
    "Restaurant": {"explanation": "short_text19nkw6pk"},
    "Breakfast room": {"explanation": "short_textibcnvl9z"},
    "Spa/Wellness": {"explanation": "short_textts4cuom6"},
    "Gym": {"explanation": "short_textzbz48iqf"},
    "Guest rooms": {"explanation": "short_text6wqdvall"},
    "Bathrooms": {"explanation": "short_textwzjcu7gg"},
    "Circulation areas (corridors, elevators, stairs)": {"explanation": "short_textc20mh2ow"},
    "Outdoor spaces/terrace": {"explanation": "short_textuutjkpb2"},
    "Meeting rooms / Event spaces": {"explanation": "short_textj7qq0acq"},
    "Back-of-house spaces (lockers, offices, staff canteen)": {"explanation": "short_textxttzxcnf"},
    "Other (General)": {"explanation": "short_textc1xaxk9u"},
    "Others space for Residential": {"explanation": "short_text4hdhx2iu"}
}


# -----------------------------
# Extract SPACE (dropdown0)
# -----------------------------
def get_selected_space_from_event(event):
    col_vals = event.get("columnValues", {})
    dd = col_vals.get("dropdown0")

    if not dd:
        return None

    chosen = dd.get("chosenValues", [])
    if not chosen:
        return None

    return chosen[0].get("name")  # Example: "Kitchen"


# -----------------------------
# Extract EXPLANATION text
# -----------------------------
def get_explanation_text(event, selected_space):
    if not selected_space:
        return ""

    col_vals = event.get("columnValues", {})
    map_entry = SPACE_EXPLANATION_MAP.get(selected_space)

    if not map_entry:
        return ""

    col_id = map_entry["explanation"]

    entry = col_vals.get(col_id, {})
    return entry.get("value", "").strip()


# ----------------------------------------
#  Convert PX → EMU
# ----------------------------------------
def px(value):
    try:
        num = float(value)
    except Exception:
        num = 0.0
    if num < 0:
        num = 0.0
    return Emu(int(round(num * 9525)))


# ----------------------------------------
#  Fit image inside box
# ----------------------------------------
def fit_center(box_w, box_h, box_x, box_y, img_w, img_h):
    if img_w <= 0 or img_h <= 0:
        return px(box_x), px(box_y), px(box_w), px(box_h)

    scale = min(box_w / img_w, box_h / img_h)
    new_w, new_h = img_w * scale, img_h * scale

    left = box_x + (box_w - new_w) / 2
    top = box_y + (box_h - new_h) / 2

    return px(left), px(top), px(new_w), px(new_h)


# ----------------------------------------
#  Get image orientation
# ----------------------------------------
def get_image_orientation(image_path):
    with Image.open(image_path) as img:
        img = ImageOps.exif_transpose(img)
        width, height = img.size

    if width > height:
        return 'H', (width, height)
    elif height > width:
        return 'V', (width, height)
    else:
        return 'S', (width, height)


# ----------------------------------------
# Inspiration layout selection
# ----------------------------------------
def choose_inspiration_slide_from_orients(orients):
    h = orients.count('H')
    v = orients.count('V')
    n = len(orients)

    if n == 1:
        return 1 if v else 2
    if n == 2:
        if h == 1 and v == 1: return 3
        if h == 2: return 4
        if v == 2: return 5
    if n == 3:
        if h == 3: return 6
        if v == 3: return 7
        if v == 2: return 8
        if h == 2: return 9
    if n == 4:
        if h == 4: return 10
        if v == 4: return 11
        if v == 3: return 12
        if h == 3: return 13
        if h == 2: return 14

    return 14


# ----------------------------------------
#  Positions for inspiration slides
# ----------------------------------------
INSPIRATION_SLIDE_BOXES = {
    1: [("V", (716, 1081, 96, 0))],
    2: [("H", (819, 538.9, 0, 540))],
    3: [("V", (491.6, 745, 101.3, 0)), ("H", (497.8, 329.7, 94.2, 757.7))],
    4: [("H", (808, 531.6, 0, 1.1)), ("H", (808, 531.6, 0, 547.2))],
    5: [("V", (630.8, 952.4, 0, 63.8)), ("V", (630.8, 952.4, 649.8, 63.8))],
    6: [("H", (1078.7, 709.7, 0, 0)), ("H", (528, 349.7, 0, 731.2)), ("H", (531.5, 349.7, 547.1, 731.2))],
    7: [("V", (711.8, 1074.6, 5.3, 5.4)), ("V", (349.7, 528, 737, 5.4)), ("V", (349.7, 528, 737, 552))],
    8: [("H", (759.9, 500, 0, 0)), ("V", (373.6, 564, 0, 516)), ("V", (374.9, 566, 389.3, 516))],
    9: [("V", (278.4, 420.3, 649.6, 660)), ("H", (928, 610.6, 0, 0)), ("H", (640, 420, 0, 660))],
    10: [("H", (1203.6, 797.2, 0, 0)), ("H", (396.3, 260.7, 1, 820.1)), ("H", (393.7, 259, 404.4, 821.9)), ("H", (396.3, 260.7, 806.5, 819.3))],
    11: [("V", (714.9, 1079.4, 20, 0)), ("V", (229.8, 346.9, 749.7, 0)), ("V", (229.8, 346.9, 749.7, 365.4)), ("V", (229.8, 346.9, 746.6, 733.1))],
    12: [("V", (301.8, 455.7, 0, 624.3)), ("V", (301.8, 455.7, 315.4, 624.3)), ("V", (301.8, 455.7, 634.9, 624.3)), ("H", (938.2, 617.3, 0, 0))],
    13: [("V", (715.3, 1080, 549.3, 0)), ("H", (528, 347.4, 0, 0)), ("H", (528, 347.4, 0, 363)), ("H", (528, 347.4, 0, 732.6))],
    14: [("H", (802.5, 528, 0, 12)), ("V", (349.7, 528, 812.1, 12)), ("V", (349.7, 528, 0, 552)), ("H", (802.5, 528, 359.3, 552))]
}


# ----------------------------------------
#  TRUE BLANK LAYOUT
# ----------------------------------------
def get_blank_layout(pres):
    for layout in pres.slide_layouts:
        if len(layout.placeholders) == 0:
            return layout
    return pres.slide_layouts[0]


# ----------------------------------------
#  MAIN SLIDE GENERATOR
# ----------------------------------------
from pptx.enum.shapes import MSO_SHAPE_TYPE

def is_picture_placeholder(shape):
    """
    Detect both PICTURE placeholders and 
    GRAPHIC_FRAME placeholders that contain an image placeholder.
    """
    NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    
    # Case 1: True placeholder flagged by python-pptx
    if shape.is_placeholder:
        return True

    # Case 2: Picture placeholder inside a graphic frame
    ph = shape._element.findall('.//{%s}ph' % NS)
    if len(ph) > 0:
        # If placeholder type is "pic" → picture placeholder
        for el in ph:
            if el.get("type") == "pic":
                return True
        return True  # treat all ph as placeholder

    return False


from PIL import Image
import os, tempfile

def prepare_image_for_box(img_path, box_w, box_h):
    """
    Fits image to box using COVER-fit:
    - No distortion
    - Scale until box filled
    - Center crop to exact box size
    """

    img = Image.open(img_path).convert("RGB")
    iw, ih = img.size

    box_ratio = box_w / box_h
    img_ratio = iw / ih

    # COVER SCALING
    if img_ratio > box_ratio:
        scale = box_h / ih    # height matches, width overflows
    else:
        scale = box_w / iw    # width matches, height overflows

    new_w = int(iw * scale)
    new_h = int(ih * scale)
    resized = img.resize((new_w, new_h), Image.LANCZOS)

    # CENTER CROP
    left = (new_w - box_w) // 2
    top = (new_h - box_h) // 2
    cropped = resized.crop((left, top, left + box_w, top + box_h))

    temp_out = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg").name
    cropped.save(temp_out, "JPEG", quality=95)

    return temp_out

def generate_inspiration_slides(template_ppt_path, inspiration_images, prs, insert_index, event):

    selected_space = get_selected_space_from_event(event)
    explanation_text = get_explanation_text(event, selected_space)

    image_paths = [p for p in inspiration_images if isinstance(p, str) and os.path.exists(p)]
    if not image_paths:
        return insert_index

    chunks = [image_paths[i:i + 4] for i in range(0, len(image_paths), 4)]
    template = Presentation(template_ppt_path)

    for group in chunks:

        # ---- detect orientation ----
        orients, sizes = [], []
        for img in group:
            orient, (w, h) = get_image_orientation(img)
            orients.append(orient)
            sizes.append((w, h))

        slide_no = choose_inspiration_slide_from_orients(orients)
        base_slide = template.slides[slide_no - 1]

        PLACEHOLDER = "Can you explain your picture selection?1"

        # ---- Replace explanation text ----
        for shape in list(base_slide.shapes):
            if shape.has_text_frame and shape.text.strip() == PLACEHOLDER:

                original_p = shape.text_frame.paragraphs[0]
                original_runs = original_p.runs
                original_font = original_runs[0].font if original_runs else None

                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                shape._element.getparent().remove(shape._element)

                tb = base_slide.shapes.add_textbox(left, top, width, height)
                tf = tb.text_frame
                tf.clear()

                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = explanation_text or ""

                if original_font:
                    run.font.name = original_font.name
                    run.font.size = original_font.size
                    run.font.bold = original_font.bold
                    run.font.italic = original_font.italic
                    if original_font.color and original_font.color.rgb:
                        run.font.color.rgb = original_font.color.rgb

                p.alignment = original_p.alignment
                p.line_spacing = original_p.line_spacing

        # ---- Create new slide ----
        blank = get_blank_layout(prs)
        new_slide = prs.slides.add_slide(blank)

        # Remove blank placeholders
        for shape in list(new_slide.shapes):
            if shape.is_placeholder:
                shape._element.getparent().remove(shape._element)

        # Maintain order
        sldIdLst = prs.slides._sldIdLst
        new_id = sldIdLst[-1]
        sldIdLst.remove(new_id)
        sldIdLst.insert(insert_index, new_id)

        # ============================================================
        # 1) COPY REAL TEMPLATE IMAGES FIRST
        #    (Skip picture placeholders to avoid them being replaced)
        # ============================================================
        for shape in base_slide.shapes:
            if (
                shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                and not is_picture_placeholder(shape)
            ):
                new_slide.shapes._spTree.append(copy.deepcopy(shape.element))

        # ============================================================
        # 2) ADD YOUR INSPIRATION IMAGES (BEHIND TEMPLATE IMAGES)
        # ============================================================
        boxes = INSPIRATION_SLIDE_BOXES[slide_no].copy()

        for i, img in enumerate(group):
            orient = orients[i]
            iw, ih = sizes[i]

            match = next(
                (
                    j for j, (bo, _) in enumerate(boxes)
                    if bo == orient or bo == "S"
                ),
                0
            )

            _, (bw, bh, bx, by) = boxes.pop(match)

            # Prepare inspiration image: cover-fit + center-crop
            cropped_img = prepare_image_for_box(img, int(bw), int(bh))

           # Insert EXACTLY into the template-defined box
            new_slide.shapes.add_picture(
                cropped_img,
                px(bx),
                px(by),
                width=px(bw),
                height=px(bh)
              )

        # ============================================================
        # 3) COPY TEMPLATE SHAPES (TEXT, LINES, DECORATIONS)
        #    → These must appear on top of everything
        # ============================================================
        for shape in base_slide.shapes:
            if shape.shape_type not in (
                MSO_SHAPE_TYPE.PICTURE,
                MSO_SHAPE_TYPE.PLACEHOLDER
            ):
                new_slide.shapes._spTree.append(copy.deepcopy(shape.element))

        insert_index += 1

    return insert_index

