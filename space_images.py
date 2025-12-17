import json
import os
import re
from typing import Dict, List
from pprint import pprint
import requests
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from inspiration_slides import generate_inspiration_slides
from image import build_slides
from floorplan_elevation_slides import  generate_floorplan_elevation_slides
from style import insert_style_slide

MONDAY_API_KEY = os.getenv("MONDAY_API_KEY", "YOUR_KEY_HERE")
MONDAY_API_URL = "https://api.monday.com/v2"


# ========================= API HELPERS ========================= #
def get_item_files(item_id: int):
    query = """
    query($item_ids: [ID!]!) {
      items(ids: $item_ids) {
        id
        assets {
          id
          name
          public_url
          file_extension
          url
        }
      }
    }
    """

    try:
        response = requests.post(
            MONDAY_API_URL,
            json={"query": query, "variables": {"item_ids": [item_id]}},
            headers={"Authorization": MONDAY_API_KEY},
            timeout=12,
        )
        data = response.json()
    except Exception as e:
        print(f"⚠️ get_item_files error: {e}")
        return []

    assets = data.get("data", {}).get("items", [{}])[0].get("assets", [])
    results = []

    for a in assets:
        url = safe_get_public_url(a)
        if url:
            results.append({
                "id": a["id"],
                "name": a["name"],
                "url": url,
                "ext": a.get("file_extension", "").lower()
            })

    return results


def safe_get_public_url(asset):
    """Get the best valid URL from an asset."""
    url = asset.get("public_url") or asset.get("url")
    return url if url and url != "null" else None


def get_file_download_url(asset_id: int):
    query = """
    query($asset_ids: [ID!]!) {
      assets(ids: $asset_ids) {
        id
        public_url
        url
      }
    }
    """

    try:
        r = requests.post(
            MONDAY_API_URL,
            json={"query": query, "variables": {"asset_ids": [asset_id]}},
            headers={"Authorization": MONDAY_API_KEY},
            timeout=10,
        )
        data = r.json()
    except Exception:
        return None

    asset = data.get("data", {}).get("assets", [{}])[0]
    return safe_get_public_url(asset)

# =================== MAPPING =================== #
SPACE_FILE_MAP = {
    "Bathroom": {
        "pictures": "fileplbovr5u",
        "floor_plans": "file1h3srtlu",
        "elevations": "file4uz9fwms",
        "inspiration": "file22mwjjxn"
    },
    "Bedroom": {
        "pictures": "file9jlbb98z",
        "floor_plans": "fileqonasnnu",
        "elevations": "file9h698ylg",
        "inspiration": "filegc3petaf"
    },
    "Dining Room": {
        "pictures": "filez8fy4euz",
        "floor_plans": "filea8f040cf",
        "elevations": "file7xkhlavk",
        "inspiration": "filervpq3oid"
    },
    "Entrance / Foyer / Hallway": {
        "pictures": "filehsmxtd29",
        "floor_plans": "filetto0n5s5",
        "elevations": "file1bb9j5zh",
        "inspiration": "file0hn4c8cj"
    },
    "Kitchen": {
        "pictures": "file1ue4pmyn",
        "floor_plans": "filedxkj5bpr",
        "elevations": "filesm6ei5sp",
        "inspiration": "file5hztykb7"
    },
    "Living Room": {
        "pictures": "file3xhh120j",
        "floor_plans": "fileb543zqov",
        "elevations": "file2msxqfsa",
        "inspiration": "filemps4yzz5"
    },
    "Outdoor area": {
        "pictures": "filenn069si7",
        "floor_plans": "filekbvd6470",
        "elevations": "filedgbi7zys",
        "inspiration": "filevp7kgxb0"
    },
    "Storage space": {
        "pictures": "file9tk9191v",
        "floor_plans": "filee63960sa",
        "elevations": "filea2s44rs4",
        "inspiration": "file56asm9v1"
    },
    "Study / Library Room": {
        "pictures": "file4mh27bw7",
        "floor_plans": "file9s65bey0",
        "elevations": "filepjhcbz16",
        "inspiration": "fileq064lpoz"
    },
    "Toilet / Powder Room": {
        "pictures": "filebuclwglz",
        "floor_plans": "fileskr2gr5b",
        "elevations": "filet5405dik",
        "inspiration": "filemuhhfxzf"
    },
    "Entertainment": {
        "pictures": "filew0rdxc4p",
        "floor_plans": "filen5clh70t",
        "elevations": "file9v6l552x",
        "inspiration": "filechtpp170"
    },
    "Lobby / Reception": {
        "pictures": "filesy1ent9a",
        "floor_plans": "filezin5hov8",
        "elevations": "filei2aqxxh9",
        "inspiration": "filevfh0mn0v"
    },
    "Bar": {
        "pictures": "filehgp5ya99",
        "floor_plans": "filelf286vtg",
        "elevations": "file5f8zrvze",
        "inspiration": "file5pkin5hn"
    },
    "Restaurant": {
        "pictures": "filebsdbxwlz",
        "floor_plans": "file5vx43zzb",
        "elevations": "filem585u1sc",
        "inspiration": "fileimhrfacl"
    },
    "Breakfast room": {
        "pictures": "filesniwn9si",
        "floor_plans": "filef3poj491",
        "elevations": "file8h4war01",
        "inspiration": "filevhrjr6iw"
    },
    "Spa / Wellness": {
        "pictures": "file1ef5w9ft",
        "floor_plans": "fileokdq82qp",
        "elevations": "fileqakni8vr",
        "inspiration": "filelndf3dcd"
    },
    "Gym": {
        "pictures": "file0rchoe30",
        "floor_plans": "filexjdydq2e",
        "elevations": "file8r13ldli",
        "inspiration": "filesxbjflrz"
    },
    "Guest rooms": {
        "pictures": "filevbs2u114",
        "floor_plans": "filehrx6bwno",
        "elevations": "file7r02b9j0",
        "inspiration": "filegqwcep9s"
    },
    "Bathrooms (Hotel)": {
        "pictures": "filesjy4yy8d",
        "floor_plans": "filejazmde8d",
        "elevations": "fileu9r5r6ai",
        "inspiration": "filetqyycu69"
    },
    "Circulation areas": {
        "pictures": "filejc0nqwn9",
        "floor_plans": "filejswo4s9c",
        "elevations": "filextbytrr0",
        "inspiration": "filetylb0vz1"
    },
    "Outdoor spaces / Terrace": {
        "pictures": "file17c1nvlg",
        "floor_plans": "fileu7ksv5iq",
        "elevations": "file9i7xfqm4",
        "inspiration": "filerd8f8jeh"
    },
    "Meeting / Event spaces": {
        "pictures": "filept2uhkgp",
        "floor_plans": "filedn79lakl",
        "elevations": "filez3dm97ms",
        "inspiration": "fileqr049w7x"
    },
    "Back-of-house spaces": {
        "pictures": "file05kq06bq",
        "floor_plans": "file9en4grlx",
        "elevations": "filefru9ndor",
        "inspiration": "file0o5rir7o"
    },
    "Other": {
        "pictures": "filevtc3rq2s",
        "floor_plans": "filevj2umvj7",
        "elevations": "file7v160ih8",
        "inspiration": "filet9kdkqbn"
    }
}


# =================== IMAGE COLLECTION LOGIC =================== #
def normalize_space_name(name: str) -> str:
    return name.strip().replace(" / ", "/").replace("/", "/")



def build_space_image_structure(event: dict, item_id: int):

    col_vals = event.get("columnValues", {})
    result = {}

    # Collect selected spaces
    selected_spaces = []

    dropdown = col_vals.get("dropdown0", {})
    if isinstance(dropdown, dict):
        selected_spaces.extend([v["name"] for v in dropdown.get("chosenValues", [])])

    multi = col_vals.get("multi_select8n5i73q8", {})
    if isinstance(multi, dict):
        selected_spaces.extend([v["name"] for v in multi.get("chosenValues", [])])

    print(f"🧩 Selected spaces: {selected_spaces}")

    # Pull fallback files once
    fallback_assets = get_item_files(item_id)

    for space in selected_spaces:

        normalized = normalize_space_name(space)
        mapping = SPACE_FILE_MAP.get(space) or SPACE_FILE_MAP.get(normalized)

        if not mapping:
            print(f"⚠️ No mapping for {space}")
            continue

        result[space] = {ftype: [] for ftype in mapping}

        for file_type, column_id in mapping.items():

            file_block = col_vals.get(column_id, {})
            files = file_block.get("files", [])

            urls = []
            for f in files:
                url = safe_get_public_url(f)
                if not url and f.get("assetId"):
                    url = get_file_download_url(f["assetId"])
                if url:
                    urls.append(url)

            if urls:
                result[space][file_type] = urls

        # fallback to item files
        if not any(result[space].values()):
            imgs = [f["url"] for f in fallback_assets if f["ext"] in ("jpg", "jpeg", "png")]
            result[space]["pictures"] = imgs

    print(f"✅ Built structure for {len(result)} spaces")
    return result



def collect_all_pictures(event: dict, item_id: int):
    structure = build_space_image_structure(event, item_id)

    urls = []
    for space, groups in structure.items():
        urls.extend(groups.get("pictures", []))

    # Remove duplicates while preserving order
    seen = set()
    clean_list = []
    for url in urls:
        if url and url not in seen:
            seen.add(url)
            clean_list.append(url)

    print(f"🖼️ Collected {len(clean_list)} unique pictures.")
    return clean_list




# =================== POWERPOINT INSERTION =================== #
def crop_center(infile, outfile, target_w, target_h):
    """Crop and resize image centered to target dimensions."""
    img = Image.open(infile)
    w, h = img.size
    target_ratio = target_w / target_h
    curr_ratio = w / h

    if curr_ratio > target_ratio:  # wide
        new_w = int(h * target_ratio)
        left = (w - new_w) // 2
        img = img.crop((left, 0, left + new_w, h))
    else:  # tall
        new_h = int(w / target_ratio)
        top = (h - new_h) // 2
        img = img.crop((0, top, w, top + new_h))

    img = img.resize((int(target_w), int(target_h)))
    img = img.convert("RGB")
    img.save(outfile, "JPEG")


def px(value):
    """Convert pixels to PowerPoint EMUs"""
    return int(value * 9525)


def insert_images_in_ppt(prs, image_urls: List[str]):
    """
    Inserts exactly 6 processed (cropped/resized) images into the PPT.
    Automatically handles PNG→JPEG conversions and all cropping logic.
    """

    if not image_urls:
        print("⚠️ No images provided to insert.")
        return prs

    # Always need exactly 6 images
    REQUIRED = 6
    image_urls = (image_urls * ((REQUIRED // len(image_urls)) + 1))[:REQUIRED]

    # Clean unique temporary names
    temp_files = [f"tmp_{i}.jpg" for i in range(REQUIRED)]
    crop_files = [f"crop_{i}.jpg" for i in range(REQUIRED)]

    # ---- Target sizes mapped by index (0–5) ----
    TARGET_SIZES = {
        0: (713.5, 1080),   # Slide 1
        1: (1117, 540),     # Slide 2
        2: (540.3, 540),    # Slide 3A
        3: (540.3, 540),    # Slide 3B
        4: (927.8, 1080),   # Slide 4
        5: (770.9, 1080)    # Slide 11
    }

    def crop_to_fill(infile, outfile, target_w, target_h):
        """Crop the image to fill the target area without distortion."""
        img = Image.open(infile)

        # Convert PNGs or indexed-color images to full RGB → FIXES 'mode P' error
        if img.mode in ("P", "RGBA"):
            img = img.convert("RGB")

        w, h = img.size
        target_ratio = target_w / target_h
        img_ratio = w / h

        # Decide crop direction
        if img_ratio > target_ratio:
            new_w = int(h * target_ratio)
            left = (w - new_w) // 2
            img = img.crop((left, 0, left + new_w, h))
        else:
            new_h = int(w / target_ratio)
            top = (h - new_h) // 2
            img = img.crop((0, top, w, top + new_h))

        img = img.resize((int(target_w), int(target_h)), Image.LANCZOS)
        img.save(outfile, "JPEG")

    # ============================================================
    #                 DOWNLOAD + CROP ALL IMAGES
    # ============================================================
    for idx, url in enumerate(image_urls):
        try:
            print(f"⬇ Downloading image {idx+1}: {url}")
            r = requests.get(url, timeout=40)
            r.raise_for_status()

            with open(temp_files[idx], "wb") as f:
                f.write(r.content)

            tw, th = TARGET_SIZES[idx]
            crop_to_fill(temp_files[idx], crop_files[idx], tw, th)

            print(f"✅ Cropped {crop_files[idx]} → {tw}×{th}")

        except Exception as e:
            print(f"⚠️ Failed to process image {idx+1}: {e}")

    print("🖼 Inserting images into PPT...")

    # ============================================================
    #                 HELPER TO INSERT BEHIND CONTENT
    # ============================================================
    def insert_picture_behind(slide, img_path, left, top, width, height):
        pic = slide.shapes.add_picture(img_path, left, top, width, height)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

    # ---- Insert into known slide locations ----
    INSERT_MAP = [
        (0, crop_files[0], 164, 0, 713.5, 1080),
        (1, crop_files[1], 803, 0, 1117, 540),
        (2, crop_files[2], 596.9, 0, 540.3, 540),
        (2, crop_files[3], 1137.2, 0, 540.3, 540),
        (3, crop_files[4], 165.4, 0, 927.8, 1080),
        (10, crop_files[5], 108, 0, 770.9, 1080),  # Slide 11
    ]

    for slide_idx, img, x, y, w, h in INSERT_MAP:
        insert_picture_behind(prs.slides[slide_idx], img, px(x), px(y), px(w), px(h))

    # ============================================================
    #                       CLEANUP
    # ============================================================
    for f in temp_files + crop_files:
        try:
            if os.path.exists(f):
                os.remove(f)
        except:
            pass

    print("🧹 Cleanup complete.")
    return prs





def safe_folder_name(name: str):
    """
    Sanitize folder name for Windows:
    - replace / and \ with _
    - strip trailing spaces
    - remove characters not allowed in Windows filenames
    """
    name = name.strip()
    name = name.replace("/", "_").replace("\\", "_")

    # remove illegal characters:  < > : " / \ | ? *
    illegal = '<>:"/\\|?*'
    for ch in illegal:
        name = name.replace(ch, "_")

    return name


def download_space_images(event: dict, item_id: int, download_root="downloaded_spaces"):
    """
    Reads selected spaces from dropdown and hotel multi-select,
    fetches all mapped file URLs according to SPACE_FILE_MAP,
    downloads the images, and stores them.
    """

    # --- Extract selected spaces ---
    col_vals = event.get("columnValues", {})
    selected_spaces = []

    # Residential selector
    if "dropdown0" in col_vals:
        dropdown_data = col_vals["dropdown0"]
        if isinstance(dropdown_data, dict) and "chosenValues" in dropdown_data:
            selected_spaces.extend(
                [v.get("name")
                 for v in dropdown_data["chosenValues"] if v.get("name")]
            )

    # ---------------------------------------------------
    # ✅ Hotel selector (YOUR REQUEST)
    # ---------------------------------------------------
    if "multi_select8n5i73q8" in col_vals:
        multi = col_vals["multi_select8n5i73q8"]
        if isinstance(multi, dict) and "chosenValues" in multi:
            selected_spaces.extend(
                [v.get("name") for v in multi["chosenValues"] if v.get("name")]
            )

    # Remove duplicates
    selected_spaces = list(dict.fromkeys(selected_spaces))

    print(f"🧩 Selected spaces: {selected_spaces}")

    # --- Fallback item files ---
    all_item_files = get_item_files(item_id)
    print(
        f"📦 Found {len(all_item_files)} total fallback files for item {item_id}")

    output = {}

    for space in selected_spaces:
        mapping = SPACE_FILE_MAP.get(space)
        if not mapping:
            print(f"⚠️ No mapping for space {space}")
            continue

        output[space] = {k: [] for k in mapping.keys()}

        for file_type, column_id in mapping.items():
            print(f"🔍 Processing {space} → {file_type}")

            block = col_vals.get(column_id)

            # elevation_measurements is TEXT not image
            if file_type == "elevation_measurements":
                text_value = block.get("text") if isinstance(
                    block, dict) else None
                if text_value:
                    print("📝 Extracted elevation measurements text")
                    output[space][file_type].append(text_value)
                else:
                    print("⚠️ No elevation_measurements text")
                continue

            # Image fields
            safe_space = safe_folder_name(space)
            save_dir = os.path.join(download_root, safe_space, file_type)
            os.makedirs(save_dir, exist_ok=True)

            if not block or not isinstance(block, dict) or "files" not in block:
                print(f"⚠️ No files in column {column_id} for {file_type}")
                continue

            for f in block["files"]:
                url = f.get("url") or f.get("public_url")
                asset_id = f.get("assetId") or f.get("id")

                if (not url or url == "null") and asset_id:
                    print(f"🔄 Fetching real download URL for asset {asset_id}")
                    url = get_file_download_url(asset_id)

                if not url or url == "null":
                    print(f"⚠️ No valid URL found for {file_type}")
                    continue

                try:
                    print(f"⬇️ Downloading: {url}")
                    r = requests.get(url, timeout=50)
                    r.raise_for_status()

                    filename = f"{asset_id or os.path.basename(url)}.jpg"
                    filepath = os.path.join(save_dir, filename)

                    with open(filepath, "wb") as fp:
                        fp.write(r.content)

                    print(f"✅ Saved: {filepath}")
                    output[space][file_type].append(filepath)

                except Exception as e:
                    print(f"❌ Failed download: {e}")

        # Fallback
        non_empty = any(
            len(files) > 0 for files in output[space].values()
            if isinstance(files, list)
        )

        if not non_empty:
            print(f"🔁 Fallback using item files for {space}")

            safe_space = safe_folder_name(space)
            fallback_dir = os.path.join(download_root, safe_space, "pictures")
            os.makedirs(fallback_dir, exist_ok=True)

            for file_info in all_item_files:
                ext = file_info.get("ext", "").lower()
                if ext in ("jpg", "jpeg", "png"):
                    url = file_info["url"]
                    try:
                        print(f"⬇️ Fallback downloading: {url}")
                        r = requests.get(url, timeout=50)
                        r.raise_for_status()

                        filename = f"{file_info['id']}.{ext}"
                        filepath = os.path.join(fallback_dir, filename)

                        with open(filepath, "wb") as fp:
                            fp.write(r.content)

                        output[space]["pictures"].append(filepath)

                    except Exception as e:
                        print(f"❌ Fallback failed: {e}")

    print("🎉 All selected spaces processed.")
    return output


def generate_layout_content_slides(template_ppt_path, layout_images, content_images, output_ppt, insert_index, event):
    """
    Apply rules for 1 layout, 2 layouts, 4–7 images, 8+ images.
    Inserts slides starting at insert_index.
    """

    # CLEAN INPUT
    def is_image(path):
        return isinstance(path, str) and os.path.exists(path)

    layout_images = [img for img in layout_images if is_image(img)]
    content_images = [img for img in content_images if is_image(img)]

    if not layout_images or not content_images:
        return insert_index

    # Split into groups of 4
    chunks = [content_images[i:i + 4]
              for i in range(0, len(content_images), 4)]
    L = len(layout_images)
    C = len(content_images)

    def call_build(layout, chunk):
        print(
            f"📄 Calling build_slides(template, layout={layout}, images={len(chunk)})")
        return build_slides(template_ppt_path, [layout], chunk, output_ppt, insert_index, event)

    # ---------- CASE 1 ----------
    if L == 1:
        print("🟦 CASE 1: Single layout")
        layout = layout_images[0]
        for chunk in chunks:
            insert_index = call_build(layout, chunk)
        return insert_index

    # ---------- CASE 2 ----------
    if L >= 2 and C >= 8:
        print("🟩 CASE 2: 2+ layouts & ≥8 images")
        for i, chunk in enumerate(chunks):
            layout = layout_images[i] if i < L else layout_images[-1]
            insert_index = call_build(layout, chunk)
        return insert_index

    # ---------- CASE 3 ----------
    if L >= 2 and 4 <= C <= 7:
        print("🟨 CASE 3: 2+ layouts & 4–7 images")
        for i, layout in enumerate(layout_images):
            chunk = chunks[i] if i < len(chunks) else chunks[-1]
            insert_index = call_build(layout, chunk)
        return insert_index

    # ---------- CASE 4 ----------
    print("🟥 CASE 4: Fallback")
    for i, chunk in enumerate(chunks):
        layout = layout_images[i] if i < L else layout_images[-1]
        insert_index = call_build(layout, chunk)

    return insert_index


def build_space_ppt(template_ppt_path, space_data, prs, selected_style, event):

    insert_index = 4
    last_floorplan_layout = None   # track only floorplans

    for space, files in space_data.items():
        print(f"\n📂 Building slides for space: {space}")

        pictures = files.get("pictures", [])
        floorplans = files.get("floor_plans", [])
        elevations = files.get("elevations", [])
        inspiration = files.get("inspiration", [])

        # ---------------------------------------------------------
        # 1️⃣ FLOORPLAN SLIDES (floorplans + pictures)
        # ---------------------------------------------------------
        if pictures and floorplans and elevations:

            # remember the layout for fallback use later
            last_floorplan_layout = floorplans[0]

            # Floorplan/Elevation combo builder (your custom module)
            insert_index = generate_floorplan_elevation_slides(
                "templates/floorplan.pptx",
                floorplans,
                pictures,
                prs,
                insert_index,
                event,
                space      # passed correctly
            )

            # Content layout slides using the actual template
            insert_index = generate_layout_content_slides(
                "templates/imageslide.pptx",
                floorplans,
                pictures,
                prs,
                insert_index,
                event
            )
            
            insert_index = generate_floorplan_elevation_slides(
                "templates/Elevation.pptx",
                elevations,
                pictures,
                prs,
                insert_index,
                event,
                space_name=space
            )
            
            template_ppt_inspiration = "templates/inspiration_slides_template1.pptx"

            insert_index = generate_inspiration_slides(
                template_ppt_inspiration,
                inspiration,
                prs,
                insert_index,
                event
            )
            
        if pictures and not floorplans and elevations:


            if last_floorplan_layout:
                print(f"✔ Using previous FLOORPLAN layout: {last_floorplan_layout}")
                insert_index = generate_floorplan_elevation_slides(
                    "templates/floorplan.pptx",
                    [last_floorplan_layout],
                    pictures,
                    prs,
                    insert_index,
                    event,
                    space
                )
                insert_index = generate_layout_content_slides(
                    "templates/imageslide.pptx",
                    [last_floorplan_layout],
                    pictures,
                    prs,
                    insert_index,
                    event
                )
            
            insert_index = generate_floorplan_elevation_slides(
                "templates/Elevation.pptx",
                elevations,
                pictures,
                prs,
                insert_index,
                event,
                space_name=space
            )
            
            template_ppt_inspiration = "templates/inspiration_slides_template1.pptx"

            insert_index = generate_inspiration_slides(
                template_ppt_inspiration,
                inspiration,
                prs,
                insert_index,
                event
            )
                     
        if pictures and floorplans and not elevations:

            # remember the layout for fallback use later
            last_floorplan_layout = floorplans[0]

            # Floorplan/Elevation combo builder (your custom module)
            insert_index = generate_floorplan_elevation_slides(
                "templates/floorplan.pptx",
                floorplans,
                pictures,
                prs,
                insert_index,
                event,
                space      # passed correctly
            )

            # Content layout slides using the actual template
            insert_index = generate_layout_content_slides(
                "templates/imageslide.pptx",
                floorplans,
                pictures,
                prs,
                insert_index,
                event
            )
            
            insert_index = generate_floorplan_elevation_slides(
                "templates/Elevation.pptx",
                elevations,
                pictures,
                prs,
                insert_index,
                event,
                space_name=space
            )
            
            template_ppt_inspiration = "templates/inspiration_slides_template1.pptx"

            insert_index = generate_inspiration_slides(
                template_ppt_inspiration,
                inspiration,
                prs,
                insert_index,
                event
            )
            
        if pictures and not floorplans and not elevations:


            if last_floorplan_layout:
                print(f"✔ Using previous FLOORPLAN layout: {last_floorplan_layout}")
                insert_index =  generate_floorplan_elevation_slides(
                    "templates/floorplan.pptx",
                    [last_floorplan_layout],
                    pictures,
                    prs,
                    insert_index,
                    event,
                    space_name=space,
                )
                insert_index = generate_layout_content_slides(
                    "templates/imageslide.pptx",
                    [last_floorplan_layout],
                    pictures,
                    prs,
                    insert_index,
                    event
                )
            
            insert_index = generate_floorplan_elevation_slides(
                "templates/Elevation.pptx",
                elevations,
                pictures,
                prs,
                insert_index,
                event,
                space_name=space
            )
            
            template_ppt_inspiration = "templates/inspiration_slides_template1.pptx"

            insert_index = generate_inspiration_slides(
                template_ppt_inspiration,
                inspiration,
                prs,
                insert_index,
                event
            )
    
    # ---------------------------------------------------------
    # 5️⃣ STYLE SLIDE
    # ---------------------------------------------------------
    if selected_style:
        insert_style_slide(prs, selected_style, insert_index)

    return prs, insert_index




