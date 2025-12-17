from pptx import Presentation
from pptx.util import Inches
from datetime import datetime, timedelta
from PIL import Image, ImageDraw
import platform
import os
import requests
import re
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os, re
import copy
from utils import cleanup_downloaded_spaces
from pptx.util import Pt
from PIL import Image
import os


import tempfile
from PIL import Image, ImageOps
from pptx.util import Emu


def px_to_emu(px):
    if not isinstance(px, (int, float)):
        raise TypeError(f"px_to_emu() expected int/float, got {px} ({type(px)})")
    return Emu(int(px * 9525))




def process_layout_slide(prs, slide, layout_img_path, bg_img_path):

    # FIXED layout values (px)
    TARGET_H_PX = 977.1     # final height ALWAYS this
    TARGET_Y_PX = 371.9     # vertical placement stays same
    TEMPLATE_WIDTH_PX = 1080  # your template width for centering

    temp_files = []

    if layout_img_path and os.path.exists(layout_img_path):

        img = Image.open(layout_img_path)
        img = ImageOps.exif_transpose(img)

        # Force vertical orientation
        if img.width > img.height:
            img = img.rotate(90, expand=True)

        w, h = img.size

        # Scale so height becomes EXACT TARGET_H_PX
        scale = TARGET_H_PX / h
        new_h = TARGET_H_PX
        new_w = w * scale

        # -----------------------------
        #  CENTER THE IMAGE HORIZONTALLY
        # -----------------------------
        center_x_px = (TEMPLATE_WIDTH_PX - new_w) / 2

        # Convert to EMU
        left = px_to_emu(center_x_px)
        top  = px_to_emu(TARGET_Y_PX)
        width_emu  = px_to_emu(new_w)
        height_emu = px_to_emu(new_h)

        # Save temp image
        tf = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
        temp_path = tf.name
        tf.close()
        temp_files.append(temp_path)

        if img.mode in ("RGBA", "LA", "P"):
            img_conv = img.convert("RGB")
        else:
            img_conv = img

        try:
            img_conv.save(temp_path, format="JPEG", quality=90)
        except Exception:
            png_tf = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
            png_path = png_tf.name
            png_tf.close()
            img_conv.save(png_path, format="PNG")
            temp_files.append(png_path)
            temp_path = png_path

        # Insert picture centered
        slide.shapes.add_picture(
            temp_path,
            left,
            top,
            width=width_emu,
            height=height_emu
        )

        print(f"✔ Centered vertical image with fixed height {TARGET_H_PX}px, width={new_w:.2f}px")

    # Background placement (unchanged)
    if bg_img_path and os.path.exists(bg_img_path):
        slide_width, slide_height = prs.slide_width, prs.slide_height
        pic = slide.shapes.add_picture(bg_img_path, 0, 0, slide_width, slide_height)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

    return slide, temp_files





# -------------------------
# Constants
# -------------------------
if platform.system() == "Windows":
    DAY_FORMAT = "%#d-%b"
else:
    DAY_FORMAT = "%-d-%b"




def px(val):
    return int(val * 9525)

# -------------------------
# Helper: Load image from URL or path
# -------------------------
def get_local_image(path_or_url, tmp_name="temp_img.png"):
    """Download image from URL or validate local path."""
    if not path_or_url:
        return None
    
    if isinstance(path_or_url, str) and path_or_url.startswith("http"):
        try:
            resp = requests.get(path_or_url, timeout=10)
            resp.raise_for_status()
            with open(tmp_name, "wb") as f:
                f.write(resp.content)
            print(f"✅ Downloaded: {tmp_name}")
            return tmp_name
        except Exception as e:
            print(f"⚠️ Failed to download image {path_or_url}: {e}")
            return None
    
    return path_or_url if os.path.exists(path_or_url) else None


# -------------------------
# Helper: Make circular crop
# -------------------------
from PIL import Image, ImageDraw
import os

def make_circle_image(img_path, output_path=None):
    """
    1. Check if image is square
    2. If not square → center-crop to the largest possible square
    3. Crop the square into a circle with transparent background
    """

    img = Image.open(img_path).convert("RGBA")
    w, h = img.size

    # -------------------------
    # 1) MAKE IMAGE SQUARE
    # -------------------------
    if w != h:
        # crop to center square
        min_dim = min(w, h)
        left   = (w - min_dim) // 2
        top    = (h - min_dim) // 2
        right  = left + min_dim
        bottom = top + min_dim

        img = img.crop((left, top, right, bottom))

    # Now the image is perfectly square
    size = img.size[0]

    # -------------------------
    # 2) CREATE CIRCLE MASK
    # -------------------------
    mask = Image.new("L", (size, size), 0)
    draw = ImageDraw.Draw(mask)
    draw.ellipse((0, 0, size, size), fill=255)

    # Apply mask
    circle_img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    circle_img.paste(img, (0, 0), mask)

    # -------------------------
    # 3) SAVE RESULT
    # -------------------------
    if not output_path:
        output_path = os.path.splitext(img_path)[0] + "_circle.png"

    circle_img.save(output_path, "PNG")
    return output_path


# -------------------------
# SIMPLE DIAGNOSTIC: Show all text in slides
# -------------------------
def show_slide_text(prs):
    """Show all text content in slides to find placeholders."""
    print("\n🔍 CHECKING ALL SLIDES FOR PLACEHOLDERS:")
    print("=" * 50)
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        print(f"\n📄 SLIDE {slide_num}:")
        
        found_placeholders = []
        for shape_idx, shape in enumerate(slide.shapes):
            try:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        print(f"  Text: '{text}'")
                        # Look for any placeholder patterns
                        if "{{" in text and "}}" in text:
                            placeholders = re.findall(r"\{\{[^}]+\}\}", text)
                            if placeholders:
                                found_placeholders.extend(placeholders)
                                print(f"    🎯 PLACEHOLDERS: {placeholders}")
            except:
                pass
        
        if not found_placeholders:
            print("  (No placeholders found)")
    
    print("=" * 50)










# ----------------------------------------
#  Project Announcement: Fit & center image
#  (Same logic – name future-proof)
# ----------------------------------------
def pa_fit_center(box_w, box_h, box_x, box_y, img_w, img_h):
    if img_w <= 0 or img_h <= 0:
        return px(box_x), px(box_y), px(box_w), px(box_h)

    scale = min(box_w / img_w, box_h / img_h)
    new_w = img_w * scale
    new_h = img_h * scale

    left = box_x + (box_w - new_w) / 2
    top = box_y + (box_h - new_h) / 2

    return px(left), px(top), px(new_w), px(new_h)





# ----------------------------------------
#  Project Announcement: Image orientation
#  (Same logic – name future-proof)
# ----------------------------------------
def pa_get_orientation(image_path):
    with Image.open(image_path) as img:
        width, height = img.size

    if width > height:
        return 'H', (width, height)
    elif height > width:
        return 'V', (width, height)
    else:
        return 'S', (width, height)




# ----------------------------------------
#  Project Announcement: Select slide layout
#  (Same logic – name future-proof)
# ----------------------------------------
def pa_select_slide_from_orients(orients):
    h = orients.count('H')
    v = orients.count('V')
    n = len(orients)

    # 1 box
    if n == 1:
        return 1 if v == 1 else 2

    # 2 boxes
    elif n == 2:
        if h == 1 and v == 1:
            return 3
        if v == 2:
            return 4
        if h == 2:
            return 5

    # 3 boxes
    elif n == 3:
        if v == 3:
            return 6
        if h == 3:
            return 7
        if h == 1 and v == 2:
            return 8
        if h == 2 and v == 1:
            return 9

    # 4 boxes
    elif n == 4:
        if h == 4:                  # ALL H
            return 10
        if v == 4:                  # ALL V
            return 11
        if v == 3 and h == 1:       # V V V H
            return 12
        if v == 2 and h == 2:       # V V H H
            return 13
        if v == 1 and h == 3:       # V H H H
            return 14

    return 14





# ----------------------------------------
#  Project Announcement: Image box positions
#  (Same data – name future-proof)
# ----------------------------------------
PA_SLIDE_BOXES = {
    1: [("V", (1084.9, 1638,-3.9, -122.4))],
    2: [("H", (1080.2, 754.2, 0, 297.4))],
    3: [("V", (589.3, 889.8, 246.1, 21.1)), ("H", (589.5, 412.2, 245.6, 916))],
    4: [("V", (507.3,771.5, 30.3,288.9)), ("V", (508.2, 767.4, 542.5, 288.9))],
    5: [("H", (863.5, 603, 108.7, 68.9)), ("H", (863.5,603, 108.7,677.4))],
    6: [("V", (324.5, 489.9, 63.3, 182.7)), ("V", (324.1, 489.3,63.2, 676.2)), ("V", (626, 983.8, 391.8, 182.7))],
    7: [("H", (611.8, 428.4, 234.6, 30.1)), ("H", (611.8, 428.4, 234.6, 460.5)), ("H", (611.8, 428.4, 234.6, 890.8))],
    8: [("H", (882.6, 617.4, 98.6, 699.7)), ("V", (440.1, 664.5, 98.6, 33.2)), ("V", (440.1, 664.5, 541.1, 33.2))],
    9: [("V", (473.7, 715.2,560, 316.9)), ("H", (509.8, 356.4, 47.7, 316.9)), ("H", (509.8, 356.4, 47.7, 676.3))],
    10: [("H", (742.4, 518.4, 169.3, 23.2)), ("H", (742.4, 518.4, 169.3, 807.7)), ("H", (366.4, 257.4, 169.3, 546.4)), ("H", (366.4, 257.4, 541.8, 546.4))],
    11: [("V", (427.1, 644.8, 111.8, 28.5)), ("V", (427.1, 644.8, 542.9, 28.5)), ("V", (427.1, 644.8, 111.8, 677.2)), ("V", (427.1, 644.8, 542.9, 677.2))],
    12: [("V", (608, 917.9,12.4, -0.4)), ("V", (444.1, 670.6, 624.6, -0.4)), ("V", (444.1, 670.6, 624.6, 675.5)), ("H", (608.6, 426.6, 12.4, 923.2))],
    13: [("H", (710.6, 496.8, 17.4, 175)), ("H", (710.6, 496.8, 17.4, 677.5)), ("V", (329.9, 498.1, 733.7, 175)), ("V", (329.9, 498.1, 733.7, 677.7))],
    14: [("V", (416.7, 647.6, 566, 660.8)), ("H", (882.6, 615.6, 98.3, 38.2)), ("H", (462, 322.2, 98.3, 660.8)), ("H", (462, 322.2, 98.3, 989))],
}



def generate_pa_slides(template_ppt_path, pa_images, prs, insert_index):
    """
    Project Announcement: Build image slides using 1–4 images per slide.
    Uses orientation-based slide selection (same logic for now).
    Inserts slides at insert_index and returns updated index.
    """

    # Clean valid image paths
    image_paths = [p for p in pa_images if isinstance(p, str) and os.path.exists(p)]
    if not image_paths:
        return insert_index

    # Split into groups of max 4 images
    chunks = [image_paths[i:i + 4] for i in range(0, len(image_paths), 4)]

    # Load template file
    template = Presentation(template_ppt_path)

    for group in chunks:

        # Detect orientation for each image
        orients, sizes = [], []
        for img in group:
            orient, (w, h) = pa_get_orientation(img)
            orients.append(orient)
            sizes.append((w, h))

        # Choose slide layout number using PA logic (same as before)
        slide_no = pa_select_slide_from_orients(orients)

        # Base slide from template
        base_slide = template.slides[slide_no - 1]

        # Add empty slide in main PPT
        blank = prs.slide_layouts[6]
        new_slide = prs.slides.add_slide(blank)

        # Move slide to desired index
        sldIdLst = prs.slides._sldIdLst
        new_id = sldIdLst[-1]
        sldIdLst.remove(new_id)
        sldIdLst.insert(insert_index, new_id)

        # Copy shapes (decorations/background) from template
        for shape in base_slide.shapes:
            element = copy.deepcopy(shape.element)
            new_slide.shapes._spTree.insert_element_before(element, "p:extLst")

        # Clone PA layout boxes
        boxes = PA_SLIDE_BOXES[slide_no].copy()

        # Insert images into their assigned boxes
        for i, img in enumerate(group):
            orient = orients[i]
            iw, ih = sizes[i]

            # Find matching orientation box
            match = next(
                (j for j, (b_orient, _) in enumerate(boxes) if b_orient == orient or b_orient == "S"),
                0
            )

            _, (bw, bh, bx, by) = boxes.pop(match)

            # Fit & center using PA version (same logic)
            left, top, width, height = pa_fit_center(bw, bh, bx, by, iw, ih)

            # Prepare image so that it fits box perfectly
            fitted_img = pa_prepare_image_for_box(img, int(bw), int(bh))

# Now add picture WITHOUT distortion 
            new_slide.shapes.add_picture(
                fitted_img,
                px(bx),
                px(by),
                width=px(bw),
                height=px(bh)
            )


        # Move index for next slide insert
        insert_index += 1

    return insert_index



from PIL import Image

def pa_prepare_image_for_box(img_path, box_w, box_h):
    """
    1. Makes image match box aspect ratio (no distortion)
    2. Does a center crop (cover mode)
    3. Returns a temp image path that is exactly box_w x box_h
    """

    img = Image.open(img_path).convert("RGB")
    iw, ih = img.size

    box_ratio = box_w / box_h
    img_ratio = iw / ih

    # --- Step 1: Fit using COVER (scale until box is filled) ---
    if img_ratio > box_ratio:
        # Image is wider → height matches box, width overflows
        scale = box_h / ih
    else:
        # Image is taller → width matches box, height overflows
        scale = box_w / iw

    new_w = int(iw * scale)
    new_h = int(ih * scale)
    resized = img.resize((new_w, new_h), Image.LANCZOS)

    # --- Step 2: Center crop to exact box size ---
    left = (new_w - box_w) // 2
    top = (new_h - box_h) // 2
    right = left + box_w
    bottom = top + box_h

    cropped = resized.crop((left, top, right, bottom))

    # --- Step 3: Save temp image ---
    temp_path = img_path + "_boxfit.jpg"
    cropped.save(temp_path, "JPEG", quality=95)

    return temp_path


# -------------------------
# Calendar helpers
# -------------------------

def build_mapping(start_date=None):
    """
    Build date mapping for calendar placeholders.
    Ensures the displayed week always starts on Monday.
    """
    if start_date is None:
        start_date = datetime.today()

    # Move start_date back to Monday of the same week
    # Monday = 0 ... Sunday = 6
    start_monday = start_date - timedelta(days=start_date.weekday())

    mapping = {}

    # Weekday headers (Mon–Sun)
    weekdays = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]
    for i, name in enumerate(weekdays, start=1):
        mapping[f"{{{{day{i}}}}}"] = name

    # Calendar day values (d1..d49)
    for i in range(1, 50):
        date = start_monday + timedelta(days=i - 1)
        mapping[f"{{{{d{i}}}}}"] = date.strftime(DAY_FORMAT)

    return mapping


def replace_text_in_frame(text_frame, mapping):
    """Replace text placeholders in a text frame."""
    for para in text_frame.paragraphs:
        for run in para.runs:
            for ph, val in mapping.items():
                if ph in run.text:
                    run.text = run.text.replace(ph, val)


def iter_shapes(shapes):
    """Recursively iterate through all shapes including groups."""
    for shp in shapes:
        yield shp
        if shp.shape_type == 6:  # group shape
            yield from iter_shapes(shp.shapes)


def update_calendar_with_bg(prs, slide, image_path, start_date=None):
    """Update calendar with background image and date mappings."""
    mapping = build_mapping(start_date)
    slide_width, slide_height = prs.slide_width, prs.slide_height

    # ---- Add background image ----
    if image_path and os.path.exists(image_path):
        pic = slide.shapes.add_picture(image_path, 0, 0,
                                       width=slide_width,
                                       height=slide_height)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
        print("✅ Calendar background added")

    # ---- Process shapes ----
    for shp in iter_shapes(slide.shapes):

        # 1️⃣ If it's a table — process ONLY table cells
        if getattr(shp, "has_table", False):
            for row in shp.table.rows:
                for cell in row.cells:
                    if hasattr(cell, "text_frame"):
                        replace_text_in_frame(cell.text_frame, mapping)

            # IMPORTANT: continue so we don't also process text_frame twice
            continue

        # 2️⃣ If it's a text frame NOT inside a table — process normally
        if getattr(shp, "has_text_frame", False):
            replace_text_in_frame(shp.text_frame, mapping)

    return slide



def replace_text_in_ppt(slide, text_map):
    """
    Format-safe text replacement.
    ✔ Keeps all font styles & sizes
    ✔ Handles split PPT runs
    ✔ Case-insensitive regex replacement
    """
    print("\n📝 Replacing text in brochure slide...")

    for shp in iter_shapes(slide.shapes):
        if not getattr(shp, "has_text_frame", False):
            continue

        tf = shp.text_frame

        for para in tf.paragraphs:

            # Collect original runs
            runs = para.runs
            if not runs:
                continue

            # 🔥 Step 1 — build combined string of paragraph text
            full_text = "".join(run.text for run in runs)

            # Step 2 — apply all replacements on full text
            new_text = full_text
            for ph, val in text_map.items():
                val_str = str(val).replace("\n", "\r")
                new_text = re.sub(ph, val_str, new_text, flags=re.IGNORECASE)

            # If nothing changed → skip rewriting
            if new_text == full_text:
                continue

            print("✔ Text updated in paragraph")

            # 🔥 Step 3 — rewrite runs WHILE preserving formatting
            idx = 0
            for run in runs:
                run_len = len(run.text)

                run.text = new_text[idx: idx + run_len]
                idx += run_len

            # If new text is longer → append extra text to last run
            if idx < len(new_text):
                runs[-1].text += new_text[idx:]


def replace_with_circle_image(slide, img_path):
    """
    Replaces {{Image1}} with a circular image placed at
    fixed coordinates:
       X = 489.4 px
       Y = 394.8 px
       W = 780 px
       H = 780 px
    """

    if not img_path or not os.path.exists(img_path):
        print(f"⚠️ Circle image not found: {img_path}")
        return slide, None

    # Create circular crop
    cropped_path = make_circle_image(img_path)

    # Remove the {{Image1}} placeholder shape
    for shape in list(slide.shapes):
        if hasattr(shape, "text") and ("{{Image1}}" in shape.text or "{{image1}}" in shape.text):
            print(f"🔄 Removing placeholder {{Image1}}")

            slide.shapes._spTree.remove(shape._element)
            break

    # ------------------------------
    # ✨ Fixed Position & Size
    # ------------------------------
    X = 489.4
    Y = 394.8
    W = 780
    H = 780

    from pptx.util import Inches

    px_to_in = lambda px: px / 96.0

    slide.shapes.add_picture(
        cropped_path,
        Inches(px_to_in(X)),
        Inches(px_to_in(Y)),
        width=Inches(px_to_in(W)),
        height=Inches(px_to_in(H))
    )

    print("✅ Circle image added at fixed position")

    return slide, cropped_path


def cleanup_temp_files(files):
    """Clean up temporary files."""
    cleaned = 0
    for f in files:
        try:
            if f and os.path.exists(f):
                os.remove(f)
                cleaned += 1
        except Exception as e:
            print(f"⚠️ Could not delete {f}: {e}")
    print(f"🧹 Cleaned up {cleaned} temporary files")

def clean_styles(raw_styles):
    """Convert style field into a clean, comma-separated string."""
    if not raw_styles:
        return ""

    # If Typeform returns as list
    if isinstance(raw_styles, list):
        return ", ".join(s.strip() for s in raw_styles if s.strip())

    # If it is a single string, normalize separators
    text = str(raw_styles).strip()

    # Replace multiple separators with commas
    for sep in [";", "|", "/", "  "]:
        text = text.replace(sep, ",")

    # Split on commas and clean up
    parts = [p.strip() for p in text.split(",") if p.strip()]

    return ", ".join(parts)


# -------------------------
# MAIN: Create Brochure (FIXED)
# -------------------------
def create_brochure_ppt(template_path, output_path, form_data,
                        circle_img, calendar_bg, layout_img, layout_bg, all_pictures):
    """
    Main function to create brochure PPT.
    Handles Slide 1 (text + circle), Slide 2 (calendar), Slide 3 (layout + background).
    Then: adds new PA-style slides for EACH space using Pictures & Inspiration only.
    """
    print(f"🔄 Creating brochure PPT from {template_path}")

    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template file not found: {template_path}")

    prs = Presentation(template_path)
    print(f"📊 Loaded presentation with {len(prs.slides)} slides")
    show_slide_text(prs)

    temp_files = []

    # -------------------------
    # Slide 1 (Text + Circle)
    # -------------------------
    circle_img = get_local_image(circle_img, "circle.png")
    if circle_img:
        temp_files.append(circle_img)

    if len(prs.slides) > 0:
        print("\n🔄 Processing Slide 1...")
        slide = prs.slides[0]

        # ---------------------------------------
        # ⭐ UNIVERSAL FIELD EXTRACTION
        # ---------------------------------------
        project_name = (
            form_data.get("Q. Project Name")
            or form_data.get("3. Your Project / Hotel Name")
            or ""
        )

        space_designed = (
            form_data.get("11. Space to be designed")
            or form_data.get("What space(s) are you looking to design")
            or ""
        )

        area = (
            form_data.get("Q. Area")
            or form_data.get("What is the area size?")
            or ""
        )

        styles_raw = (
            form_data.get("Which style's do you like")
            or form_data.get("Which style's do you like1")
            or ""
        )

        styles = clean_styles(styles_raw)

        city = form_data.get("City") or form_data.get("city") or ""
        country = form_data.get("Country1") or form_data.get("country1") or ""
        location = f"{city}, {country}".strip(", ")

        # ---------------------------------------
        # ⭐ DETECT PROJECT TYPE (Hotel/Residential)
        # ---------------------------------------
        project_type_raw = (
            form_data.get("project_type")
            or form_data.get("Project Type")
            or ""
        ).strip().lower()

        project_type_text = "Hotel" if "hotel" in project_type_raw else "Residential"

        print("➡️ Detected Project Type:", project_type_text)

        # ---------------------------------------
        # ⭐ TEXT MAP (Regex-based)
        # ---------------------------------------
        text_map = {
            r"Q\. Project Name": project_name,
            r"11\. Space to be designed": space_designed,
            r"Q\. Area": area,
            r"Which style's do you like": styles,
            r"Location1": location,

            # 🔥 Replace ANY version of “Residential”
            r"(?i)residential": project_type_text,
        }

        replace_text_in_ppt(slide, text_map)

        if circle_img:
            slide, cropped_circle = replace_with_circle_image(slide, circle_img)
            if cropped_circle:
                temp_files.append(cropped_circle)

    # -------------------------
    # Slide 2 (Calendar)
    # -------------------------
    calendar_bg = get_local_image(calendar_bg, "calendar_bg.png")
    if calendar_bg:
        temp_files.append(calendar_bg)
        if len(prs.slides) > 1:
            slide = prs.slides[1]
            update_calendar_with_bg(prs, slide, calendar_bg)

    # -------------------------
    # Slide 3 (Layout)
    # -------------------------
    layout_img = get_local_image(layout_img, "layout_img.png")
    layout_bg = get_local_image(layout_bg, "layout_bg.png")

    if layout_img:
        temp_files.append(layout_img)
    if layout_bg:
        temp_files.append(layout_bg)

    if len(prs.slides) > 2:
        slide = prs.slides[2]
        process_layout_slide(prs, slide, layout_img, layout_bg)

    # -------------------------
    # Space-based PA slides
    # -------------------------
    insert_at = 3  

    for space, files in all_pictures.items():
        print(f"\n➡️ Space: {space}")

        categories = {
            "Pictures": files.get("pictures", []),
            "Inspiration": files.get("inspiration", []),
        }

        for category_name, images in categories.items():
            if not images:
                continue

            valid_images = [img for img in images if img and os.path.exists(img)]
            if not valid_images:
                continue

            chunks = [valid_images[i:i + 4] for i in range(0, len(valid_images), 4)]

            for chunk in chunks:
                insert_at = generate_pa_slides(
                    template_ppt_path="templates/pa_slides.pptx",
                    pa_images=chunk,
                    prs=prs,
                    insert_index=insert_at
                )

    # -------------------------
    # SAVE OUTPUT
    # -------------------------
    print(f"\n💾 Saving presentation to: {output_path}")
    prs.save(output_path)
    print(f"✅ Brochure PPT created: {output_path}")

    cleanup_temp_files(temp_files)
    cleanup_downloaded_spaces(base_path="downloaded_spaces", exclude_folder=None)

    return output_path






FILES_DIR = "files"
os.makedirs(FILES_DIR, exist_ok=True)


def generate_brochure(item_id, selected_styles, form_data, all_pictures, event):
    """Generate the brochure PPT file."""
    BTEMPLATE_PATH = "templates/brochure.pptx"
    BOUTPUT_PATH   = os.path.join(FILES_DIR, f"{item_id}_boutput.pptx")

    # ---------------------------------------------------
    # ⭐ Detect Project Type (Hotel / Residential)
    # ---------------------------------------------------
    group_name = (event.get("groupName") or "").lower()

    if "hotel" in group_name:
        project_type = "Hotel"
    else:
        project_type = "Residential"

    print("\n🏷️ DETECTED PROJECT TYPE:", project_type)

    # -----------------------------
    # Helper: Safe wrapped indexing
    # -----------------------------
    def get_wrapped(img_list, index, fallback=None):
        """
        Return img_list[index], but if index is out of range,
        wrap around using modulo. If img_list empty → return fallback.
        """
        if not img_list:
            return fallback
        return img_list[index % len(img_list)]

    # ---------------------------------------------------
    # 1. Collect FIRST IMAGES from each category per space
    # ---------------------------------------------------
    first_pictures   = []
    first_floorplans = []

    for space, files in all_pictures.items():

        print(f"\n📂 Processing space: {space}")

        pictures    = files.get("pictures", [])
        floorplans  = files.get("floor_plans", [])

        # collect first image for pictures
        if pictures:
            first_pictures.append(pictures[0])

        # collect first image for floorplans
        if floorplans:
            first_floorplans.append(floorplans[0])

    # ---------------------------------------------------
    # 2. Assign theme images using wrapped fallback logic
    # ---------------------------------------------------
    circle_img  = get_wrapped(first_pictures,   0, None)
    calendar_bg = get_wrapped(first_pictures,   1, circle_img)
    layout_img  = get_wrapped(first_floorplans, 0, circle_img)
    layout_bg   = get_wrapped(first_pictures,   2, layout_img)

    print("\n🖼️ Selected Theme Images:")
    print("circle_img: ", circle_img)
    print("calendar_bg:", calendar_bg)
    print("layout_img: ", layout_img)
    print("layout_bg:  ", layout_bg)

    # ---------------------------------------------------
    # 3. Pass EVERYTHING (including project_type) to PPT builder
    # ---------------------------------------------------
    return create_brochure_ppt(
        BTEMPLATE_PATH,
        BOUTPUT_PATH,
        form_data={**form_data, "project_type": project_type},  # 👈 injected here
        circle_img=circle_img,
        calendar_bg=calendar_bg,
        layout_img=layout_img,
        layout_bg=layout_bg,
        all_pictures=all_pictures,
    )
