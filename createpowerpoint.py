from email_utils import send_email_with_ppt
from pptx import Presentation
import os
import re
import json
import requests
from city import insert_city_image_in_ppt
from space_images import collect_all_pictures, insert_images_in_ppt, download_space_images, build_space_ppt
from Hotel import collect_all_hotel_pictures, insert_images_hotel_ppt, build_hotel_space_ppt, map_hotel_form_data , insert_logo_hotel_ppt
from image import  SPACE_INFO


FILES_DIR = "files"
os.makedirs(FILES_DIR, exist_ok=True)
TEMPLATE_PATH = "template.pptx"
BASE_URL = os.getenv("BASE_URL")

MONDAY_API_KEY = os.getenv("MONDAY_API_KEY", "eyJhbGciOiJIUzI1NiJ9.eyJ0aWQiOjU0NjI5MjM1NywiYWFpIjoxMSwidWlkIjo3NDc3Njk5NywiaWFkIjoiMjAyNS0wOC0wNFQwOTo0MzowNS4wMDBaIiwicGVyIjoibWU6d3JpdGUiLCJhY3RpZCI6MTIxNDMyMDQsInJnbiI6InVzZTEifQ.yYeelRXHOZlaxwYHBAvi6eXRzD2fNn1H-jX-Pd8Ukcw")
MONDAY_API_URL = "https://api.monday.com/v2"

# Track processed items to prevent duplicate emails
PROCESSED_ITEMS = set()
EMAIL_SENT_LOG = "email_sent_log.txt"


def load_processed_items():
    """Load previously processed item IDs from file"""
    global PROCESSED_ITEMS
    if os.path.exists(EMAIL_SENT_LOG):
        try:
            with open(EMAIL_SENT_LOG, 'r') as f:
                PROCESSED_ITEMS = set(line.strip()
                                      for line in f if line.strip())
            print(
                f"📋 Loaded {len(PROCESSED_ITEMS)} previously processed items")
        except Exception as e:
            print(f"⚠️ Error loading processed items: {e}")
            PROCESSED_ITEMS = set()
    else:
        PROCESSED_ITEMS = set()


def mark_item_as_processed(item_id):
    """Mark an item as processed and save to file"""
    global PROCESSED_ITEMS
    PROCESSED_ITEMS.add(str(item_id))
    try:
        with open(EMAIL_SENT_LOG, 'a') as f:
            f.write(f"{item_id}\n")
        print(f"✅ Marked item {item_id} as processed")
    except Exception as e:
        print(f"⚠️ Error saving processed item: {e}")


def is_item_processed(item_id):
    """Check if an item has already been processed"""
    return str(item_id) in PROCESSED_ITEMS


# Load processed items on startup
load_processed_items()



def calculate_image_size_for_slide_fixed_height(img_width, img_height, placeholder_width, placeholder_height, fixed_height_px=8000):
    """
    Calculate image size with FIXED height of 540px and proportional width.

    Args:
        img_width, img_height: Original image dimensions in pixels
        placeholder_width, placeholder_height: Placeholder dimensions in PowerPoint EMUs
        fixed_height_px: Fixed height in pixels (default 540)

    Returns:
        (new_width, new_height) in PowerPoint units (EMUs)
    """
    print(f"🔍 Input: img={img_width}x{img_height}px, placeholder={placeholder_width}x{placeholder_height}EMUs, fixed_height={fixed_height_px}px")

    # Convert pixels to PowerPoint EMUs (English Metric Units)
    # 1 inch = 914400 EMUs, assuming 72 DPI (standard for most images)
    PIXELS_PER_INCH = 72
    EMUS_PER_INCH = 914400

    # Calculate aspect ratio (width/height)
    aspect_ratio = img_width / img_height if img_height > 0 else 1.0
    print(f"📐 Original aspect ratio (w/h): {aspect_ratio:.3f}")

    # Set fixed height and calculate proportional width
    new_height_px = fixed_height_px
    new_width_px = int(fixed_height_px * aspect_ratio)

    print(f"📏 Fixed dimensions: {new_width_px}x{new_height_px}px")

    # Convert to EMUs
    new_width_emu = int((new_width_px / PIXELS_PER_INCH) * EMUS_PER_INCH)
    new_height_emu = int((new_height_px / PIXELS_PER_INCH) * EMUS_PER_INCH)

    # Check if image fits within placeholder bounds
    if new_width_emu > placeholder_width:
        scale_factor = placeholder_width / new_width_emu
        new_width_emu = placeholder_width
        new_height_emu = int(new_height_emu * scale_factor)
        final_height_px = int(new_height_emu / EMUS_PER_INCH * PIXELS_PER_INCH)
        print(
            f"⚠️ Width exceeded placeholder, scaled down to: {new_width_px}x{final_height_px}px")

    print(f"📏 Final EMU dimensions: {new_width_emu}x{new_height_emu}")
    print(
        f"📊 Final pixel equivalent: {new_width_emu/EMUS_PER_INCH*PIXELS_PER_INCH:.0f}x{new_height_emu/EMUS_PER_INCH*PIXELS_PER_INCH:.0f}px")

    return new_width_emu, new_height_emu


def get_file_download_url(asset_id: int) -> str:
    """Get the actual downloadable S3 URL for a Monday.com file using the API"""
    query = """
    query($asset_ids: [ID!]!) {
      assets(ids: $asset_ids) {
        id
        name
        public_url
        file_extension
        url
      }
    }
    """
    headers = {
        "Authorization": MONDAY_API_KEY,
        "Content-Type": "application/json"
    }

    try:
        response = requests.post(
            MONDAY_API_URL,
            json={
                "query": query,
                "variables": {"asset_ids": [str(asset_id)]}
            },
            headers=headers,
            timeout=15
        )
        response.raise_for_status()
        data = response.json()

        if "errors" in data:
            print(f"⚠️ GraphQL errors for asset {asset_id}: {data['errors']}")
            return None

        assets = data.get("data", {}).get("assets", [])
        if assets:
            asset = assets[0]
            public_url = asset.get("public_url") or asset.get("url")
            if public_url and public_url != "null":
                return public_url

        return None

    except Exception as e:
        print(f"⚠️ Failed to get URL for asset {asset_id}: {e}")
        return None

# PDF / DOCX / ZIP Image Extraction




def compute_combined_area(col_values):
    """
    Reads selected spaces from dropdown0,
    fetches their area fields, and returns TOTAL area:
        e.g., '80 m2'
    """
    selected = col_values.get("dropdown0", {}).get("chosenValues", [])
    selected_spaces = [item["name"] for item in selected]

    if not selected_spaces:
        return ""

    total_area = 0

    for space in selected_spaces:
        if space not in SPACE_INFO:
            continue

        area_col_id = SPACE_INFO[space]["total_area"]
        area_col = col_values.get(area_col_id, {})

        raw_val = area_col.get("value") or area_col.get("text") or ""

        # Ensure numeric
        if raw_val.strip().isdigit():
            total_area += int(raw_val)

    if total_area == 0:
        return ""

    return f"{total_area} "
def clean_styles(raw_styles):
    """Convert style field into a clean, comma-separated string."""
    if not raw_styles:
        return ""

    # If Typeform returns as list
    if isinstance(raw_styles, list):
        return ", ".join(s.strip() for s in raw_styles if s.strip())

    # If it is a single string, normalize separators
    text = str(raw_styles).strip()

    # Replace multiple separators with commas
    for sep in [";", "|", "  "]:
        text = text.replace(sep, ",")

    # Split on commas and clean up
    parts = [p.strip() for p in text.split(",") if p.strip()]

    return ", ".join(parts)



def map_residential_form_data(event: dict):
    col = event.get("columnValues", {})

    # Extract image URLs from Monday "files" type column
    image_urls = []
    if "files" in col and col["files"].get("value"):
        try:
            file_data = json.loads(col["files"]["value"])
            image_urls = [asset["url"] for asset in file_data]
        except Exception:
            image_urls = []

    # Extract style information from dropdown column
    styles = []
    style_columns = ["dropdown", "dropdown0", "dropdown1", "dropdown2", "style_dropdown"]

    for col_name in style_columns:
        if col_name in col and col[col_name].get("chosenValues"):
            chosen_values = col[col_name].get("chosenValues", [])
            styles = [v.get("name", "") for v in chosen_values if v.get("name")]
            if styles:
                print(f"🎨 Found styles in column '{col_name}': {styles}")
                break

    styles_text = ",".join(styles) if styles else ""
    styles = clean_styles(styles_text)

    # ✅ NEW: compute combined area using helper
    combined_area = compute_combined_area(col)

    return {
        "9. What is the property type": col.get("dropdown76", {}).get("chosenValues", [{}])[0].get("name"),
        "City": col.get("text8", {}).get("value"),
        "city": col.get("text8", {}).get("value"),
        "Country1": col.get("country6", {}).get("countryName"),
        "country1": col.get("country6", {}).get("countryName"),
        "11. Space to be designed": ", ".join([v.get("name", "") for v in col.get("dropdown0", {}).get("chosenValues", [])]),
        "What is the area size?": col.get("short_text8fr4spel", {}).get("value"),
        "Which style's do you like1": styles,
        "5. How old are you1": col.get("status", {}).get("label", {}).get("text"),
        "12. How many people will leave in the space1": col.get("text1", {}).get("value"),
        "10. What best describes your situation1": col.get("single_selecti4d0sw1", {}).get("label", {}).get("text"),
        "13. Kids1": col.get("text2", {}).get("value"),
        "14. Do you have any pets1": col.get("text_1", {}).get("value"),
        "16. Please describe the scope of work": col.get("text37", {}).get("value"),
        "22. Is there any other information…": col.get("long_text3", {}).get("text"),
        "Can you explain your picture selection?1": col.get("short_textot656d98", {}).get("value"),
        "15. What words describe best the mood and feel": col.get("short_text5fonuzuu", {}).get("value"),
      

        # ✅ NEW FIELD ADDED HERE
        "Q. Area": combined_area,
    }







def fetch_user_details(email: str):
    try:
        url = f"https://migrate.omrsolutions.com/get_user_details.php?email={email}"
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            return response.json()
    except Exception as e:
        print("Error fetching user details:", e)
    return {}


def get_item_files(item_id: int):
    """Get all files for an item with their public URLs"""
    query = """
    query($item_ids: [ID!]!) {
      items(ids: $item_ids) {
        id
        name
        assets {
          id
          name
          public_url
          file_extension
          url
        }
      }
    }
    """
    headers = {"Authorization": MONDAY_API_KEY,
               "Content-Type": "application/json"}
    try:
        response = requests.post(
            MONDAY_API_URL,
            json={
                "query": query,
                "variables": {"item_ids": [str(item_id)]}
            },
            headers=headers,
            timeout=15
        )
        response.raise_for_status()
        data = response.json()

        if "errors" in data:
            print(f"⚠️ GraphQL errors for item {item_id}: {data['errors']}")
            return []

    except Exception as e:
        print(f"⚠️ Request/JSON error: {e}")
        return []

    items = data.get("data", {}).get("items", [])
    if not items:
        return []

    assets = items[0].get("assets", [])
    result = []

    for asset in assets:
        public_url = asset.get("public_url") or asset.get("url")
        if public_url and public_url != "null":
            result.append({
                "id": asset["id"],
                "name": asset["name"],
                "url": public_url,
                "ext": asset["file_extension"]
            })

    return result


def replace_text_in_ppt(prs, text_map: dict):
    """Enhanced text replacement with multiline support and better debugging.
       Same logic as your original function, but modifies the existing prs object.
    """

    print("🔄 Replacing text inside Presentation object...")

    replacements_made = 0
    all_found_text = []  # To debug what placeholders exist

    # --- First pass: collect all text ---
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue

            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                for r_idx, run in enumerate(p.runs):
                    if run.text.strip():
                        all_found_text.append(
                            f"Slide {slide_idx+1}: '{run.text.strip()}'"
                        )

    print("📝 All text found in template (first 20 entries):")
    for text in all_found_text[:20]:
        print(f"  {text}")

    # --- Second pass: perform replacements ---
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue

            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                for r_idx, run in enumerate(p.runs):

                    for placeholder, value in text_map.items():
                        if not value:
                            continue

                        # Convert \n to PowerPoint line breaks
                        value_str = str(value).replace("\n", "\r")

                        # Exact match
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, value_str)
                            replacements_made += 1
                            print(
                                f"✅ Exact match: '{placeholder}' -> '{value_str}' in slide {slide_idx + 1}"
                            )

                        # Case-insensitive match
                        elif placeholder.lower() in run.text.lower():
                            pattern = re.compile(
                                re.escape(placeholder), re.IGNORECASE)
                            new_text = pattern.sub(value_str, run.text)

                            if new_text != run.text:
                                run.text = new_text
                                replacements_made += 1
                                print(
                                    f"✅ Case-insensitive: '{placeholder}' -> '{value_str}' in slide {slide_idx + 1}"
                                )

    print(f"📝 Total replacements made: {replacements_made}")

    # --- Debug unused placeholders ---
    print("\n🔍 Form data that might not have been used:")
    for key, value in text_map.items():
        if value and key not in ''.join(all_found_text):
            print(f"  '{key}': '{value}'")

    return prs


def categorize_and_collect_images(event: dict) -> dict:
    """
    Categorize images from Monday.com files based on column source.
    - First tries to fetch a public_url (S3).
    - If no public_url, downloads via Monday API and extracts images locally.
    - Only processes JPG, JPEG, and PNG files - skips all other file types
    """
    categorized_images = {
        "floor_plans": [],
        "elevation_drawings": [],
        "existing_pictures": [],
        "inspiration_images": []
    }

    col_vals = event.get("columnValues", {})
    print(f"🔍 Found {len(col_vals)} column values to process")

    # Map file columns to categories
    column_category_map = {
        "files": "floor_plans",
        "fileb3p8t108": "elevation_drawings",
        "fileh7us51cr": "existing_pictures",
        "files3": "inspiration_images"
    }

    # Supported image extensions
    SUPPORTED_EXTENSIONS = ["jpg", "jpeg", "png"]

    for column_name, category in column_category_map.items():
        print(f"🔍 Checking column: {column_name} -> {category}")

        if column_name not in col_vals:
            print(f"⚠️ Column {column_name} not found in webhook data")
            continue

        file_data = col_vals[column_name]
        if not (isinstance(file_data, dict) and "files" in file_data):
            print(f"⚠️ Unexpected structure in {column_name}: {file_data}")
            continue

        files_list = file_data["files"]
        print(f"📁 Found {len(files_list)} files in {column_name}")

        for file_info in files_list:
            if not isinstance(file_info, dict):
                print(f"⚠️ File info is not a dict: {file_info}")
                continue

            asset_id = file_info.get("assetId") or file_info.get("id")
            filename = file_info.get("name", "")
            file_ext = file_info.get("extension", "").lower()

            if not asset_id or not filename:
                print(f"⚠️ Missing asset_id or filename: {file_info}")
                continue

            # ✅ CHECK: Skip non-image files
            if file_ext not in SUPPORTED_EXTENSIONS:
                print(
                    f"⏭️ SKIPPING: {filename} (extension: {file_ext}) - Only JPG, JPEG, PNG supported")
                continue

            print(
                f"🔍 Processing image: {filename} (ID: {asset_id}, Ext: {file_ext})")

            try:
                # 1. Try to get a direct public URL
                public_url = get_file_download_url(asset_id)

                if public_url and ("amazonaws.com" in public_url or "s3" in public_url):
                    categorized_images[category].append(public_url)
                    print(f"✅ Public URL added to {category}: {public_url}")
                    continue

                # 2. If no public URL, download via Monday API
                print(
                    f"⬇️ Downloading image {filename} from Monday API (no public_url)")
                headers = {"Authorization": MONDAY_API_KEY}
                download_url = f"{MONDAY_API_URL}/file/{asset_id}"
                response = requests.get(
                    download_url, headers=headers, timeout=20)
                response.raise_for_status()
                file_bytes = response.content

                # 3. Save image as temporary file
                temp_path = f"temp_{asset_id}.{file_ext}"
                with open(temp_path, "wb") as f:
                    f.write(file_bytes)

                categorized_images[category].append(temp_path)
                print(f"✅ Image saved and added to {category}: {temp_path}")

            except Exception as e:
                print(f"⚠️ Failed to process {filename}: {e}")
                import traceback
                traceback.print_exc()

    # Remove empty categories
    categorized_images = {k: v for k, v in categorized_images.items() if v}

    print("📁 Image categorization complete:")
    for category, urls in categorized_images.items():
        print(f"  {category}: {len(urls)} images")
        for i, url in enumerate(urls, 1):
            print(f"    {i}. {url}")

    return categorized_images


def filter_style_slides_optimized(prs, selected_styles):
    """Filter style slides based on selected styles - placeholder function"""
    # This function should be implemented based on your style filtering logic
    print(f"🎨 Filtering slides for styles: {selected_styles}")
    return prs


def count_page(prs, placeholder="NO."):
    print("📄 Page numbering started.")

    counter = 1

    for slide_index, slide in enumerate(prs.slides, start=1):
        replaced = False

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:

                    if placeholder in run.text:
                        # Add slash + leading zeros + dots
                        page_str = f"/{counter:02d}"

                        run.text = run.text.replace(placeholder, page_str)
                        replaced = True

        if replaced:
            print(f"✅ Slide {slide_index}: Inserted page number {page_str}")
        else:
            print(f"⚠️ Slide {slide_index}: No placeholder found")

        counter += 1

    print("📄 Page numbering complete.")
    return prs








# ============================================================
#                HELPER FUNCTIONS
# ============================================================

def extract_event(body):
    """Extract Monday event or challenge handshake."""
    if "challenge" in body:
        return {"challenge": body["challenge"]}, True
    if "event" not in body:
        return None, False
    return body["event"], False


def check_if_processed(event):
    """Avoid processing same item twice."""
    item_id = event.get("pulseId")
    if is_item_processed(item_id):
        return True, {
            "status": "skipped",
            "message": f"Item {item_id} already processed",
            "item_id": item_id
        }
    return False, None


import re

def normalize_style_list(styles_str):
    if not styles_str:
        return []

    # Remove quotes
    cleaned = styles_str.replace('"', "").replace("'", "")

    # Replace " and " with comma
    cleaned = re.sub(r"\s+and\s+", ",", cleaned, flags=re.IGNORECASE)

    # Special rule: "Asian, Zen" → "Asian Zen"
    cleaned = re.sub(r"\bAsian\s*,\s*Zen\b", "Asian Zen", cleaned, flags=re.IGNORECASE)

    # Also catch "Asian / Zen"
    cleaned = re.sub(r"\bAsian\s*[/\\]\s*Zen\b", "Asian Zen", cleaned, flags=re.IGNORECASE)

    # Replace line breaks & semicolons
    cleaned = re.sub(r"[\n;]+", ",", cleaned)

    # Normalize repeated spaces
    cleaned = re.sub(r"\s+", " ", cleaned).strip()

    parts = [p.strip() for p in cleaned.split(",") if p.strip()]

    # Deduplicate while keeping order
    seen = set()
    result = []
    for p in parts:
        key = p.lower()
        if key not in seen:
            seen.add(key)
            result.append(p)

    return result


def extract_form_data(event):
    """Extract mapped form data, Monday columns, and selected styles."""

    group_name = (event.get("groupName") or "").lower()

    # ---------------------------------------------------
    # 🗂️ USE DIFFERENT FORM MAPPERS BASED ON PROJECT TYPE
    # ---------------------------------------------------
    if "residentiel" in group_name:
        print("🟦 Using RESIDENTIAL form mapper...")
        # ← your residential mapping function
        form_data = map_residential_form_data(event)
    elif "hotel" in group_name:
        print("🟩 Using HOTEL form mapper...")
        # ← your hotel mapping function
        form_data = map_hotel_form_data(event)

    # ---------------------------------------------
    # Monday raw column values
    # ---------------------------------------------
    col_vals = event.get("columnValues", {})

    # Potential style keys depending on form type
    possible_style_keys = [
        "Which style's do you like1",
        "27. Style",
        "Which styles do you like",
        "Styles",
    ]

    styles_str = ""

    # Check all possible keys
    for key in possible_style_keys:
        if key in form_data and form_data[key]:
            styles_str = form_data[key]
            break

    # Normalize: "Coastal, Classic" OR "Coastal\nClassic"
    selected_styles = normalize_style_list(styles_str)


    return form_data, col_vals, selected_styles


def extract_email(col_vals, form_data):
    """Extract email from Monday submission or fallback."""
    email = None

    if "email" in col_vals:
        email_block = col_vals["email"]
        if isinstance(email_block, dict):
            email = email_block.get("email") or email_block.get("text")

    if not email:
        email = (
            form_data.get("Email")
            or form_data.get("email")
           
        )
    return email


def update_form_with_db(form_data, email):
    """Fetch DB details and merge into form_data."""
    user_details = fetch_user_details(email)

    if user_details.get("status") != "success":
        return form_data, None

    qd = user_details["data"].get("quotationdetails", {})

    if qd.get("area_size"):
        form_data["Q. Area"] = qd["area_size"]

    if qd.get("project_name"):
        form_data["Q. Project Name"] = qd["project_name"]

    if qd.get("residential_type"):
        form_data["Q.Nature of the project"] = qd["residential_type"]

    return form_data, qd


def resolve_project_name(form_data, col_vals, event):
    """Ensure project name always has final fallback."""
    if form_data.get("Q. Project Name"):
        return form_data

    if event.get("pulseName"):
        form_data["Q. Project Name"] = event["pulseName"]
        return form_data

    name = (col_vals.get("text30", {}) or {}).get("value")
    form_data["Q. Project Name"] = name or "Unknown Project"

    return form_data


def generate_main_ppt(event, item_id, selected_styles, form_data, project_type):
    group_name = (event.get("groupName") or "").lower()
    print("this is event", event)

    print("🔎 Event group name:", group_name)
    print("📌 project_type parameter:", project_type)

    # Normalize project type
    project_type = project_type.lower().strip()

    # Supported groups
    allowed_groups = ["residentiel projects", "hotel projects"]

    # ❌ Skip unsupported groups
    if group_name not in allowed_groups:
        print(f"❌ PPT generation skipped. Unsupported group: {group_name}")
        return None

    # -----------------------------------------------------------
    # 📄 SELECT TEMPLATE BASED ON PROJECT TYPE
    # -----------------------------------------------------------
    if "residentiel" in group_name:
        TEMPLATE_PATH = "templates/residential_template.pptx"
    elif "hotel" in group_name:
        TEMPLATE_PATH = "templates/hotel_template.pptx"
    else:
        TEMPLATE_PATH = "template.pptx"   # fallback (should never happen)

    OUTPUT_PATH = os.path.join(FILES_DIR, f"{item_id}_output.pptx")

    print(f"📄 Using template: {TEMPLATE_PATH}")

    # Load base PPT
    prs = Presentation(TEMPLATE_PATH)

    # -----------------------------------------------------------
    # 🏠 RESIDENTIAL PROJECTS
    # -----------------------------------------------------------
    if "residentiel" in group_name:
        print("🏠 Generating Residential Project PPT...")

        prs = insert_city_image_in_ppt(prs, form_data.get("city"))

        all_pictures = collect_all_pictures(event, item_id)
        insert_images_in_ppt(prs, all_pictures)

        space_data = download_space_images(event, item_id)
        build_space_ppt("templates/imageslide.pptx", space_data, prs ,selected_styles , event)
        
      

        replace_text_in_ppt(prs, form_data)
        
        count_page(prs, placeholder="/NO.")
    

    # -----------------------------------------------------------
    # 🏨 HOTEL PROJECTS
    # -----------------------------------------------------------
    elif "hotel" in group_name:
        print("🏨 Generating Hotel Project PPT...")

    
        prs = insert_city_image_in_ppt(prs, form_data.get("city"))
        
        hotel_pictures = collect_all_hotel_pictures(event, item_id)
        insert_images_hotel_ppt(prs, hotel_pictures)
        insert_logo_hotel_ppt(prs, event)

        hotel_space = download_space_images(event, item_id)
        build_hotel_space_ppt("templates/imageslide.pptx", hotel_space, prs,selected_styles,event)
        
        count_page(prs, placeholder="NO.")
        replace_text_in_ppt(prs, form_data)
        

    # -----------------------------------------------------------
    # Save file
    # -----------------------------------------------------------
    prs.save(OUTPUT_PATH)
    print(f"✅ PPT saved: {OUTPUT_PATH}")

    return OUTPUT_PATH


def extract_project_type(item):
    """
    Finds the project type (Residential / Hotel / F&B / Other)
    from the column with id 'status1'.
    """
    for col in item.get("column_values", []):
        if col.get("id") == "status1":
            return (col.get("text") or "").strip()

    return ""  # default if nothing found


def send_final_email(item_id, form_data, email):
    """Send final email with PPT download links."""
    
    # Use hardcoded email from .env instead of form email
    recipient_email = os.getenv("SEND_TO_EMAIL")
    if not recipient_email:
        print("⚠️ SEND_TO_EMAIL not set in .env file, using form email as fallback")
        recipient_email = email
    
    output_url = f"{BASE_URL}/download-ppt?item_id={item_id}&ppt_type=output"
    brochure_url = f"{BASE_URL}/download-ppt?item_id={item_id}&ppt_type=brochure"

    project_name = form_data.get("Q. Project Name", "Your Project")

    subject = f"Your Project Files Are Ready – {project_name}"

    html_content = f"""
    <div style="max-width:600px;margin:auto;font-family:Arial,Helvetica,sans-serif;color:#333;">
        
        <!-- Header -->
        <div style="background:#F5C518;padding:20px;text-align:center;border-radius:10px 10px 0 0;">
            <h2 style="margin:0;color:#1A1A1A;font-weight:900;">
                THE LIVE DESIGN PROJECT
            </h2>
        </div>

        <!-- Body -->
        <div style="padding:25px;background:#ffffff;border:1px solid #eee;border-top:0;">
            <p>Hello 👋,</p>

            <p>Your project <b>{project_name}</b> is ready!</p>

            <p>You can download your presentation files below:</p>

            <div style="margin:20px 0;">
                <a href="{output_url}"
                   style="display:inline-block;background:#F5C518;color:#000;
                          padding:12px 22px;border-radius:25px;text-decoration:none;
                          font-weight:bold;margin:5px 0;">
                    Download Project Brief
                </a>
            </div>

            <div style="margin:20px 0;">
                <a href="{brochure_url}"
                   style="display:inline-block;background:#F5C518;color:#000;
                          padding:12px 22px;border-radius:25px;text-decoration:none;
                          font-weight:bold;margin:5px 0;">
                    Download Project Announcement
                </a>
            </div>

            <p style="margin-top:30px;">
                If you have any questions or want to discuss your project further,
                feel free to reply to this email.
            </p>

            <p>Warm regards,<br>The Live Design Project Team</p>
        </div>

        <!-- Footer -->
        <div style="text-align:center;font-size:12px;color:#777;margin-top:15px;">
            No pressure. Just good design advice from friendly professionals.
        </div>
    </div>
    """

    try:
        send_email_with_ppt(
            recipient=recipient_email,
            subject=subject,
            html_content=html_content,
            sender_email=os.getenv("BREVO_SENDER_EMAIL"),
            sender_name="The Live Design Project",
            ppt_paths=[]
        )
        print(f"✅ Email sent successfully to {recipient_email}")

    except Exception as e:
        print(f"⚠️ Failed to send email: {e}")
        import traceback
        traceback.print_exc()
