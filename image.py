from pptx import Presentation
from pptx.util import Emu
from PIL import Image
import os
import copy
from PIL import Image, ImageOps  

from pptx.slide import Slide
import copy
import os
from PIL import Image, ImageOps
import tempfile

# ------------------------
# 🔹 Helper: Convert pixel → EMU
def px(value):
    """Convert pixel to EMU (PowerPoint unit)."""
    try:
        num = float(value)
    except Exception:
        num = 0.0
    if num < 0:
        num = 0.0
    return Emu(int(round(num * 9525)))


# ------------------------
# 🔹 Fit and center image while preserving aspect ratio
def fit_center(box_w, box_h, box_x, box_y, img_w, img_h):
    if img_w <= 0 or img_h <= 0:
        return px(box_x), px(box_y), px(box_w), px(box_h)
    scale = min(box_w / img_w, box_h / img_h)
    new_w, new_h = img_w * scale, img_h * scale
    left = box_x + (box_w - new_w) / 2
    top = box_y + (box_h - new_h) / 2
    return px(left), px(top), px(new_w), px(new_h)


from PIL import Image, ImageOps

def get_image_orientation(image_path):
    with Image.open(image_path) as img:

        # FIX: correct orientation using EXIF metadata
        img = ImageOps.exif_transpose(img)

        width, height = img.size

    if width > height:
        return 'H', (width, height)
    elif height > width:
        return 'V', (width, height)
    else:
        return 'S', (width, height)


# ------------------------
# 🔹 Choose correct slide layout
def choose_ppt_slide_from_orients(orients):
    h = orients.count('H')
    v = orients.count('V')
    n = len(orients)

    if n == 1:
        return 1 if v == 1 else 2
    elif n == 2:
        if h == 1 and v == 1: return 3
        if h == 2: return 4
        if v == 2: return 5
    elif n == 3:
        if h == 3: return 6
        if v == 3: return 7
        if v == 2 and h == 1: return 8
        if v == 1 and h == 2: return 9
    elif n == 4:
        if h == 4: return 10
        if v == 4: return 11
        if v == 3 and h == 1: return 12   # 3V + 1H
        if h == 3 and v == 1: return 13   # 3H + 1V
        if h == 2 and v == 2: return 14
    raise ValueError("Combination not supported (1–4 images only).")





SLIDE_BOXES = {
    1: [("V", (716, 1081, 96, 0))],
    2: [("H", (819, 538.9, 0, 540))],
    3: [("V", (716, 1081, 101.3, 0)), ("H", (635.3, 420.8, 832.8, 660.3))],
    4: [("H", (808, 531.6, 0, 1.1)), ("H", (808, 531.6, 0, 547.2))],
    5: [("V", (415.9, 628, 0, 0)), ("V", (415.9, 628, 432.1, 0))],
    6: [("H", (626.3, 414.8, 0, 666.1)), ("H", (626.3, 414.8, 641.9, 666.1)), ("H", (626.3, 414, 1287.6, 666.1))],
    7: [("V", (559.5, 844.7, 5.3, 5.4)), ("V", (349.7, 528, 565, 0.6)), ("V", (349.7, 528, 564, 556))],
    8: [("H", (760, 500, 1.7, 0)), ("V", (373.6, 564, 1.7, 516)), ("V", (374.9, 566, 389.3, 516))],
    9: [("V", (326.5, 492.9, 66.5, 587.5)), ("H", (747.8, 492, 409.6, 588)), ("H", (747.8, 492, 1171.4, 588))],
    10: [("H", (633.2, 419.4, -0.3, 231.4)), ("H", (633.2, 416.6, 0, 664.3)), ("H", (633.2, 416.6, 644.7, 664.3)), ("H", (632.7, 416.3, 1286.8, 664.3))],
    11: [("V", (355.4, 536.5, 201.9, 0)), ("V", (355.4, 536.5, 567.8, 0)), ("V", (355.4, 536.5, 203.4, 545.7)), ("V", (355.4, 536.5, 564.8, 545.7))],
    12: [("V", (319.2, 481.9, 201.1, 597.6)), ("V", (319.2, 481.9, 531.4, 597.6)), ("V", (319.2, 481.9, 859.5, 597.6)), ("H", (732.4, 481.9, 1187.6, 597.6))],
    13: [("V", (436.8, 659.4, 386.7, 0)), ("H", (632, 415.8, 0, 665.1)), ("H", (632, 415.8, 644, 665.1)), ("H", (632, 415.8, 1288, 665.1))],
    14: [("H", (886.3, 583.2, 1.7, 0)), ("V", (320.8, 484.4, 1.7, 595.6)), ("V", (320.8, 484.4, 328, 595.6)), ("H", (733.2, 484.4, 656, 595.6))],
}

LAYOUT_SLIDE_BOXES = {
    1: [(1095.5, 659.3, 824, 0)],
    2: [(1095.5, 659.3, 824, 0)],
    3: [(1095.5, 659.3, 824, 0)],
    4: [(1095.5, 659.3, 824, 0)],
    5: [(1055.4, 627.4, 864.1, 0)],
    6: [(1095.5, 659.3, 824, 0)],
    7: [(995.1, 587.4, 924.3, 0)],
    8: [(1134.8, 691.4, 784.7, 0)],
    9: [(977.3, 580.6, 942.2, 0)],
    10: [(1079.5, 650.2, 840, 0)],
    11: [(995.1, 587.4, 924.3, 0)],
    12: [(995.1, 587.4, 924.3, 0)],
    13: [(1087.5, 659.4, 832, 0)],
    14: [(1015.5, 582.6, 904, 0)],
}

SPACE_INFO = {
    "Bathroom": {
        "elevation_measurements": "long_texti11yq90s",
        "total_area": "short_textxryyyyn1"
    },
    "Bedroom": {
        "elevation_measurements": "long_text3fq8nrj3",
        "total_area": "short_text8ufxux4l"
    },
    "Dining Room": {
        "elevation_measurements": "long_texttonzaaop",
        "total_area": "short_textcje3k6c6"
    },
    "Entrance / Foyer / Hallway": {
        "elevation_measurements": "long_text08nejzfw",
        "total_area": "short_textel0sxqar"
    },
    "Kitchen": {
        "elevation_measurements": "long_textvbhtdci5",
        "total_area": "short_textvlbq48x0"
    },
    "Living Room": {
        "elevation_measurements": "long_text9rkb5lmt",
        "total_area": "short_textdke9knnw"
    },
    "Outdoor area": {
        "elevation_measurements": "long_textj9cu0t6c",
        "total_area": "short_text3n6cra8h"
    },
    "Storage space": {
        "elevation_measurements": "long_texto26evbeb",
        "total_area": "short_textvl5bpmfg"
    },
    "Study / Library Room": {
        "elevation_measurements": "long_textn91xo3pc",
        "total_area": "short_textnd4qfrsb"
    },
    "Toilet / Powder Room": {
        "elevation_measurements": "long_textshycswxa",
        "total_area": "short_text267d0qe8"
    },
    "Entertainment": {
        "elevation_measurements": "long_textkd6v8vo2",
        "total_area": "short_text82yr051d"
    },

    # -----------------------------
    # Hotel / Commercial Spaces
    # -----------------------------
    "Lobby": {
        "elevation_measurements": "short_textlbwha8xg",
        "total_area": "short_text6jjx97mb"
    },
    "Bar": {
        "elevation_measurements": "short_textdi0idew7",
        "total_area": "short_textk92m8cfh"
    },
    "Restaurant": {
        "elevation_measurements": "short_text1t5cdy0s",
        "total_area": "short_text6lq5zjiz"
    },
    "Breakfast Room": {
        "elevation_measurements": "short_textqai6m1xo",
        "total_area": "short_textqkig0f4q"
    },
    "Spa/Wellness": {
        "elevation_measurements": "short_textc5pxfard",
        "total_area": "short_textripyd7ok"
    },
    "Gym": {
        "elevation_measurements": "short_text3hz110lk",
        "total_area": "short_textyf9cgnwi"
    },
    "Guest rooms": {
        "elevation_measurements": "short_textwfzf4yev",
        "total_area": "short_text8l1bxdy7"
    },
    "Bathrooms": {
        "elevation_measurements": "short_text4qcd13u0",
        "total_area": "short_textvz2zbwzp"
    },
    "Circulation areas": {
        "elevation_measurements": "short_textsk1upzgu",
        "total_area": "short_textvilnaa9l"
    },
    "Outdoor spaces": {
        "elevation_measurements": "short_textyl093m0x",
        "total_area": "short_textk2hwmq1i"
    },
    "Meeting rooms / Event spaces": {
        "elevation_measurements": "short_text4p5xg50n",
        "total_area": "short_textscfu64dq"
    },
    "Back-of-house spaces": {
        "elevation_measurements": "short_textuo9xzeif",
        "total_area": "short_text7b477vir"
    },
    "Other": {
        "elevation_measurements": "short_textauyl3mp5",
        "total_area": "short_textxudlnksd"
    }
}


def fit_image_in_layout_box(layout_box, img_path):
    """
    Fits and centers image in a layout box, adjusting for orientation differences.
    layout_box: (width, height, x, y)
    """
    bw, bh, bx, by = layout_box
    orient, (iw, ih) = get_image_orientation(img_path)

    # if image is vertical but box is horizontal, we can try fitting by height primarily
    box_ratio = bw / bh
    img_ratio = iw / ih

    # determine scale based on which dimension dominates
    if orient == 'V' and img_ratio < box_ratio:
        # Vertical image in a horizontal box -> fit by height
        scale = bh / ih
    else:
        # Normal fit by smallest dimension
        scale = min(bw / iw, bh / ih)

    new_w, new_h = iw * scale, ih * scale
    left = bx + (bw - new_w) / 2
    top = by + (bh - new_h) / 2

    return px(left), px(top), px(new_w), px(new_h)


def prepare_image_for_box(img_path, box_w, box_h):
    """
    Fits image to box using COVER-fit:
    - No distortion
    - Scale until box filled
    - Center crop to exact box size
    """

    img = Image.open(img_path).convert("RGB")
    iw, ih = img.size

    box_ratio = box_w / box_h
    img_ratio = iw / ih

    # COVER SCALING
    if img_ratio > box_ratio:
        scale = box_h / ih    # height matches, width overflows
    else:
        scale = box_w / iw    # width matches, height overflows

    new_w = int(iw * scale)
    new_h = int(ih * scale)
    resized = img.resize((new_w, new_h), Image.LANCZOS)

    # CENTER CROP
    left = (new_w - box_w) // 2
    top = (new_h - box_h) // 2
    cropped = resized.crop((left, top, left + box_w, top + box_h))

    temp_out = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg").name
    cropped.save(temp_out, "JPEG", quality=95)

    return temp_out


def add_images_to_ppt(template_ppt_path, image_paths, layout, output_ppt, insert_index, event):
    
    template = Presentation(template_ppt_path)

    # -----------------------------
    # ORIENTATION
    # -----------------------------
    orients, sizes = [], []
    for img_path in image_paths:
        orient, size = get_image_orientation(img_path)
        orients.append(orient)
        sizes.append(size)

    slide_no = choose_ppt_slide_from_orients(orients)

    # 🔥 DEBUG — show exactly which template is being loaded
    import os
    print("\n========== TEMPLATE DEBUG ==========")
    print("template_ppt_path:", template_ppt_path)
    print("ABSOLUTE PATH:", os.path.abspath(template_ppt_path))
    print("Actual template slide count:", len(template.slides))
    print("Requested slide_no:", slide_no)
    print("====================================\n")

    # 🔥 SAFETY CHECK — prevents crash if template has fewer slides
    if slide_no < 1 or slide_no > len(template.slides):
        raise ValueError(
            f"ERROR: Template has only {len(template.slides)} slides, "
            f"but choose_ppt_slide_from_orients returned slide_no={slide_no}. "
            f"Fix your template or SLIDE_BOXES definitions."
        )

    # SAFE NOW — no crash here
    base_slide = template.slides[slide_no - 1]

    # -----------------------------
    # CREATE NEW SLIDE
    # -----------------------------
    blank_layout = get_blank_layout(output_ppt)
    new_slide = output_ppt.slides.add_slide(blank_layout)

    # remove any placeholders that may exist
    for shape in list(new_slide.shapes):
        if shape.is_placeholder:
            el = shape._element
            el.getparent().remove(el)

    # reorder to correct insert_index
    sldIdLst = output_ppt.slides._sldIdLst
    new_slide_id = sldIdLst[-1]
    sldIdLst.remove(new_slide_id)
    sldIdLst.insert(insert_index, new_slide_id)

    # -----------------------------
    # COPY TEMPLATE SHAPES
    # -----------------------------
    for shape in base_slide.shapes:
        element = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(element, "p:extLst")

    # ================================================================
    # ⭐ EMPTY ELEVATION HANDLING — NO LAYOUT IMAGE
    # ================================================================
    if layout == "EMPTY_ELEVATION":
        print("🚫 No Elevation Layout Image — leaving layout box empty.")
    else:
        print("\n🖼️ Processing Main Layout Image...")

        img = Image.open(layout)
        img = ImageOps.exif_transpose(img)

        layout_w, layout_h = img.size

        if layout_h > layout_w:
            print("🔄 Layout image is VERTICAL — Rotating to HORIZONTAL")
            img = img.rotate(90, expand=True)

        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
        img.convert("RGB").save(temp_file.name)
        layout_final = temp_file.name

        iw, ih = img.size
        bw, bh, bx, by = LAYOUT_SLIDE_BOXES[slide_no][0]
        left, top, width, height = fit_center(bw, bh, bx, by, iw, ih)

        new_slide.shapes.add_picture(layout_final, left, top, width=width, height=height)

    # -----------------------------
    # INSERT CONTENT IMAGES
    # -----------------------------
    boxes = SLIDE_BOXES[slide_no].copy()

    for i, img_path in enumerate(image_paths):
        orient = orients[i]
        iw, ih = sizes[i]

        match_idx = next(
            (j for j, (bo, _) in enumerate(boxes) if bo == orient or bo == "S"),
            0
        )

        _, (bw, bh, bx, by) = boxes.pop(match_idx)

        # Prepare inspiration image: cover-fit + center-crop
        cropped_img = prepare_image_for_box(img_path, int(bw), int(bh))


        # Insert EXACTLY into the template-defined box
        new_slide.shapes.add_picture(
            cropped_img,
            px(bx),
            px(by),
            width=px(bw),
           height=px(bh)
           )


    return insert_index + 1




# ------------------------
# 🔹 Build Slides
def build_slides(template_ppt_path, layout_images, content_images, output_ppt, insert_index,event):

    print("\n🚀 Starting slide generation...")
    
    if not layout_images:
        raise ValueError("Layout image required.")
    if not content_images:
        raise ValueError("Content images required.")

    groups = [content_images[i:i + 4] for i in range(0, len(content_images), 4)]

    for idx, image_group in enumerate(groups):
        layout = layout_images[idx % len(layout_images)]
        print(f"\n==============================================")
        print(f"📄 Creating slide {idx+1} using layout {layout}")
        print("==============================================")

        insert_index = add_images_to_ppt(
            template_ppt_path, image_group, layout,
            output_ppt, insert_index, event
        )

    print("\n🎉 All slides created successfully!")
    return insert_index



def get_blank_layout(pres):
    """Return a layout that has no placeholders (true blank)."""
    for layout in pres.slide_layouts:
        if len(layout.placeholders) == 0:
            return layout
    return pres.slide_layouts[0]


