from pptx import Presentation
from pptx.util import Inches
import requests
import os
from PIL import Image
from io import BytesIO

def save_image(url, city_name, width=489, height=540):
    headers = {"User-Agent": "CityImageFetcher/1.0"}
    try:
        img_response = requests.get(url, headers=headers, stream=True, verify=False)
        if img_response.status_code == 200:
            # Open image with PIL
            img = Image.open(BytesIO(img_response.content))
            # Resize
            img = img.resize((width, height), Image.Resampling.LANCZOS)

            filename = f"{city_name.lower()}.jpg"
            img.save(filename, "JPEG")
            print(f"✅ Image of {city_name} saved as {filename} with size {width}x{height}")
            return filename
        else:
            print(f"❌ Failed to download {city_name} image. Status: {img_response.status_code}")
            return None
    except Exception as e:
        print(f"❌ Error downloading {city_name}: {e}")
        return None


def get_city_image(city_name):
    url = "https://en.wikipedia.org/w/api.php"
    params = {
        "action": "query",
        "titles": city_name,
        "prop": "pageimages",
        "format": "json",
        "pithumbsize": 600
    }
    headers = {"User-Agent": "CityImageFetcher/1.0"}

    response = requests.get(url, params=params, headers=headers).json()
    pages = response.get("query", {}).get("pages", {})

    for _, page in pages.items():
        if "thumbnail" in page:
            img_url = page["thumbnail"]["source"]
            print(f"✅ Found Wikipedia thumbnail for {city_name}: {img_url}")
            return save_image(img_url, city_name)

    print(f"❌ No image found for {city_name}.")
    return None


def insert_city_image_in_ppt(prs, city_name):
    """Insert city image into slide 2 at a FIXED position.
       Ignores placeholder size and uses the coordinates you provided.
    """

    # Fetch and save temporary city image
    img_file = get_city_image(city_name)
    if not img_file:
        print("⚠ No image to insert in PPT")
        return prs

    # Access slide 2 (index 1)
    try:
        slide = prs.slides[1]
    except:
        print("⚠ Slide 2 not found in PPT")
        return prs

    # ---------------------------------------------
    # 🔍 Remove placeholder {{LocationImage}}
    # ---------------------------------------------
    placeholder_found = False
    for shape in list(slide.shapes):
        if shape.has_text_frame and "{{LocationImage}}" in shape.text:
            slide.shapes._spTree.remove(shape._element)
            placeholder_found = True
            print("🗑️ Removed {{LocationImage}} placeholder")
            break

    if not placeholder_found:
        print("⚠ No {{LocationImage}} placeholder found on slide 2 (image will still be added)")

    # ---------------------------------------------
    # 📌 FIXED IMAGE POSITION
    # ---------------------------------------------
    X = 313.9
    Y = 0
    W = 489.1
    H = 540

    from pptx.util import Inches
    px_to_in = lambda px: px / 96.0

    slide.shapes.add_picture(
        img_file,
        Inches(px_to_in(X)),
        Inches(px_to_in(Y)),
        width=Inches(px_to_in(W)),
        height=Inches(px_to_in(H))
    )

    print(f"✅ Inserted city image for '{city_name}' at fixed coordinates")

    # ---------------------------------------------
    # 🗑️ Cleanup
    # ---------------------------------------------
    try:
        os.remove(img_file)
        print(f"🗑️ Deleted temporary image file: {img_file}")
    except:
        pass

    return prs
