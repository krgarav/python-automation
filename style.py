from pptx import Presentation

STYLE_TO_SLIDE_MAP = {
    "Asian/Zen": -1,
    "Coastal": -2,
    "Art Deco": -3,
    "Contemporary": -4,
    "Country": -5,
    "Eclectic": -6,
    "Industrial": -7,
    "Mid-century": -8,
    "Modern": -9,
    "Minimalist": -10,
    "Rustic": -11,
    "Scandinavian": -12,
    "Transitional": -13,
    "Tropical": -14,
    "Urban": -15,
    "Shabby Chic": -16,
    "Traditional": -17,
}

def get_slide_index_by_style(prs, style_name):
    """Return absolute slide index for built-in style slide."""
    if style_name not in STYLE_TO_SLIDE_MAP:
        return None

    neg_index = STYLE_TO_SLIDE_MAP[style_name]   # e.g. -1, -2, -3
    index = len(prs.slides) + neg_index          # convert to absolute index

    if index < 0 or index >= len(prs.slides):
        return None

    return index


def move_slide(prs, old_index, new_index):
    """Moves slide without corrupting _sldIdLst."""
    sldIdLst = prs.slides._sldIdLst

    # extract slide node
    slide_id = sldIdLst[old_index]

    # remove old
    del sldIdLst[old_index]

    # adjust new index if slide was removed before it
    if new_index > old_index:
        new_index -= 1

    sldIdLst.insert(new_index, slide_id)


def delete_slide(prs, index):
    """Safely delete slide by index."""
    sldIdLst = prs.slides._sldIdLst
    del sldIdLst[index]

def remove_leftover_style_slides(prs, used_style_count):
    TOTAL_STYLE_SLIDES = 17
    slides_to_remove = TOTAL_STYLE_SLIDES - used_style_count

    print(f"Removing leftover style slides: {slides_to_remove}")

    for _ in range(slides_to_remove):
        last_index = len(prs.slides) - 1
        delete_slide(prs, last_index)
        print(f"✔ Removed leftover slide at index {last_index}")
def insert_style_slide(prs, styles, insert_index):


    # Normalize to list
    if isinstance(styles, str):
        styles = [styles]

    # --- SORT STYLES BY THEIR MAP VALUE (correct order) ---
    styles = sorted(styles, key=lambda s: STYLE_TO_SLIDE_MAP.get(s, 0))

    print(f"Sorted styles (by mapping): {styles}")

    insert_index = insert_index

    # Move slides in mapping order
    for style_name in styles:
        old_index = get_slide_index_by_style(prs, style_name)
        if old_index is None:
            print(f"⚠ Style '{style_name}' not found, skipping…")
            continue

        move_slide(prs, old_index, insert_index)
        print(f"✔ Moved '{style_name}' from {old_index} → {insert_index}")
        insert_index += 1

    # Remove leftover unused slides
    remove_leftover_style_slides(prs, len(styles))

    return insert_index