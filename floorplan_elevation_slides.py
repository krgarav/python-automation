from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from PIL import Image, ImageOps
import copy
import io
import os
from image import SPACE_INFO

PX = 9525  # EMUs per “pixel-ish” unit


# ---------------------------------------------------------
# BASIC HELPERS
# ---------------------------------------------------------

def get_blank_layout(prs: Presentation):
    """Return a blank layout (no placeholders), or first layout as fallback."""
    for layout in prs.slide_layouts:
        if len(layout.placeholders) == 0:
            return layout
    return prs.slide_layouts[0]


def is_unwanted_template_image(shape) -> bool:
    """Skip tiny/logo/preview images from template."""
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return False

    name = shape.name.lower()
    if "thumb" in name or "thumbnail" in name or "preview" in name or "logo" in name:
        return True

    w, h = shape.width, shape.height
    if w < 760000 and h < 760000:
        return True

    return False


# ---------------------------------------------------------
# SAFE SHAPE COPY
# ---------------------------------------------------------

def safe_copy_shape(src_shape, dst_slide):
    st = src_shape.shape_type

    # Copy pictures safely
    if st == MSO_SHAPE_TYPE.PICTURE:
        img = src_shape.image
        blob = img.blob
        stream = io.BytesIO(blob)

        new_pic = dst_slide.shapes.add_picture(
            stream,
            left=src_shape.left,
            top=src_shape.top,
            width=src_shape.width,
            height=src_shape.height,
        )
        new_pic.rotation = src_shape.rotation
        return

    # All other shapes → deep-copy XML append
    new_el = copy.deepcopy(src_shape.element)
    dst_slide.shapes._spTree.append(new_el)


# ---------------------------------------------------------
# IMAGE INSERT HELPERS
# ---------------------------------------------------------



def add_main_simple_image(slide, image_path: str):
    """Insert image preserving aspect ratio and cropping instead of stretching."""

    frame_width = int(879.3 * PX)
    frame_height = int(540 * PX)
    left = int(0 * PX)
    top = 0

    # Insert at natural size
    picture = slide.shapes.add_picture(
        image_path,
        left=left,
        top=top,
    )

    # Aspect ratios
    img_ratio = picture.width / picture.height
    frame_ratio = frame_width / frame_height

    if img_ratio > frame_ratio:
        # Crop left & right
        extra = img_ratio/frame_ratio - 1
        picture.crop_left = extra/2
        picture.crop_right = extra/2
    else:
        # Crop top & bottom
        extra = frame_ratio/img_ratio - 1
        picture.crop_top = extra/2
        picture.crop_bottom = extra/2

    # Resize to frame using integer values
    picture.width = frame_width
    picture.height = frame_height
    picture.left = left
    picture.top = top

    return picture

def add_layout_full_height_on_right(slide, image_path: str, prs: Presentation):
    """Insert the tall right-side layout image."""
    SLIDE_WIDTH = prs.slide_width
    PX_LOCAL = 9525
    target_height_px = 1080

    img = Image.open(image_path)
    img = ImageOps.exif_transpose(img)

    if img.width > img.height:
        img = img.rotate(90, expand=True)

    img = img.convert("RGB")
    scale = target_height_px / img.height
    final_h = target_height_px * PX_LOCAL
    final_w = img.width * scale * PX_LOCAL

    temp = image_path + "_layout.jpg"
    img.save(temp, "JPEG")

    left = SLIDE_WIDTH - final_w

    return slide.shapes.add_picture(
        temp,
        left=int(left),
        top=0,
        width=int(final_w),
        height=int(final_h),
    )





# ---------------------------------------------------------
# TEMPLATE SLIDE LOADER
# ---------------------------------------------------------

def get_template_first_slide(template_ppt_path: str):
    """Always load a fresh presentation."""
    prs = Presentation(template_ppt_path)
    return prs.slides[0]


# ---------------------------------------------------------
# BUILD ONE ELEVATION SLIDE
# ---------------------------------------------------------

def build_elevation_slide(template_ppt_path: str,
                          pictures,
                          elevation_img: str,
                          output_ppt: Presentation,
                          insert_index: int,
                          space_name: str,
                          event: dict):
    """
    Build a single elevation slide, insert images, copy template shapes,
    and replace measurement text based on space_name.
    """

    # Load a fresh template slide
    template_slide = get_template_first_slide(template_ppt_path)

    # Create blank slide
    blank_layout = get_blank_layout(output_ppt)
    slide = output_ppt.slides.add_slide(blank_layout)

    # Remove placeholders on new slide
    for shp in list(slide.shapes):
        if shp.is_placeholder:
            el = shp._element
            el.getparent().remove(el)

    # Insert slide at provided index
    sldIdLst = output_ppt.slides._sldIdLst
    new_slide_id = sldIdLst[-1]
    sldIdLst.remove(new_slide_id)
    sldIdLst.insert(insert_index, new_slide_id)

    # -----------------------------------------------------
    # 1️⃣ ADD IMAGES FIRST (behind template)
    # -----------------------------------------------------
    if pictures:
       best_img = get_best_horizontal_image(pictures)
       if best_img:
          add_main_simple_image(slide, best_img)

    if elevation_img != "EMPTY_ELEVATION":
        add_layout_full_height_on_right(slide, elevation_img, output_ppt)

    # -----------------------------------------------------
    # 2️⃣ COPY TEMPLATE SHAPES ON TOP
    # -----------------------------------------------------
    for shp in template_slide.shapes:
        if is_unwanted_template_image(shp):
            continue
        safe_copy_shape(shp, slide)

    # -----------------------------------------------------
    # 3️⃣ REPLACE MEASUREMENTS TEXT IN THE SLIDE
    # -----------------------------------------------------
    replace_measurements_in_slide(slide, space_name, SPACE_INFO, event)

    return insert_index + 1

def resolve_monday_value(field_id: str, event: dict):
    cols = event.get("column_values") or event.get("columnValues") or {}
    
    for key, col in cols.items():
        if key == field_id:
            return col.get("text") or col.get("value") or ""
    return ""

def normalize_text(value: str) -> str:
    """Normalize PPT text: fix weird spaces, remove hidden chars, collapse spacing."""
    if not value:
        return ""

    # Fix PPT non-breaking spaces
    value = value.replace("\xa0", " ")

    # Remove invisible characters
    for hidden in ["\u200b", "\u200c", "\u200d"]:
        value = value.replace(hidden, "")

    # Collapse multiple spaces
    value = " ".join(value.split())

    return value


def replace_measurements_in_slide(slide, space_name: str, SPACE_INFO: dict, event: dict):
    print("\n---------------- DEBUG: replace_measurements_in_slide ----------------")

    # 1) Pull space info
    space_data = SPACE_INFO.get(space_name)
    print("SPACE NAME:", space_name)
    print("SPACE_DATA:", space_data)

    if not space_data:
        print("ERROR: No space data found.")
        return

    # 2) Extract Monday IDs
    elev_id = space_data.get("elevation_measurements")
    print("Elevation Field ID:", elev_id)

    # 3) Get real values from Monday event + normalize
    raw = resolve_monday_value(elev_id, event) if elev_id else ""
    print("Raw Monday Value:", repr(raw))

    elev_text = normalize_text(raw)
    print("Normalized Value (elev_text):", repr(elev_text))

    # 4) Replacement rules
    REPLACE_MAP = {
        "XXXXX": elev_text,
    }
    print("REPLACE_MAP:", REPLACE_MAP)

    print("\n--- LOOP SHAPES ---")
    # 5) Loop all shapes
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:

            merged = "".join(run.text for run in paragraph.runs)
            print("\nOriginal paragraph text:", repr(merged))

            norm = normalize_text(merged)
            print("Normalized paragraph text:", repr(norm))

            # Apply replacements
            for key, val in REPLACE_MAP.items():
                print(f"Trying to replace '{key}' with '{val}'")
                if val:
                    norm = norm.replace(key, val)
                else:
                    print(f"SKIPPED: Value for {key} is empty.")

            print("Final paragraph text:", repr(norm))

            # Write back
            if paragraph.runs:
                paragraph.runs[0].text = norm
                for r in paragraph.runs[1:]:
                    r.text = ""

    print("---------------- END DEBUG ----------------\n")


# ---------------------------------------------------------
# PUBLIC ENTRY
# ---------------------------------------------------------

def generate_floorplan_elevation_slides(template_ppt_path: str,
                                        elevation_images,
                                        pictures,
                                        output_ppt: Presentation,
                                        insert_index: int,
                                        event: dict,
                                        space_name: str):

    elevation_items = elevation_images if elevation_images else [
        "EMPTY_ELEVATION"]

    for elevation in elevation_items:
        insert_index = build_elevation_slide(
            template_ppt_path,
            pictures,
            elevation,
            output_ppt,
            insert_index,
            space_name,
            event
        )

    return insert_index





def get_best_horizontal_image(picture_list):
    """
    Returns a horizontal image path.
    - If any picture is already horizontal → return it.
    - If all are vertical → crop the first vertical image to horizontal (center crop).
    """

    if not picture_list:
        return None

    horizontal = []
    vertical = []

    # --- Separate images by orientation ---
    for img_path in picture_list:
        try:
            img = Image.open(img_path)
            img = ImageOps.exif_transpose(img)

            if img.width >= img.height:
                horizontal.append(img_path)
            else:
                vertical.append(img_path)

        except Exception as e:
            print(f"⚠️ Failed to inspect image {img_path}: {e}")
            continue

    # --- CASE 1: We already have a horizontal image ---
    if horizontal:
        return horizontal[0]

    # --- CASE 2: All images vertical → crop one ---
    if vertical:
        src_path = vertical[0]
        img = Image.open(src_path)
        img = ImageOps.exif_transpose(img)

        # Horizontal center crop (no blank padding)
        new_width = img.height
        left = (img.width - new_width) // 2
        right = left + new_width
        cropped = img.crop((left, 0, right, img.height))

        # Ensure compatibility with JPEG
        if cropped.mode in ("P", "RGBA"):
            cropped = cropped.convert("RGB")

        # Generate a safe temp filename without pathlib
        base, ext = os.path.splitext(src_path)
        temp_path = f"{base}_hcrop.jpg"

        cropped.save(temp_path, "JPEG")
        return temp_path

    return None
