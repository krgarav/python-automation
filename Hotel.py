import json
import os
import re
from typing import Dict, List
import requests
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from inspiration_slides import generate_inspiration_slides
from space_images import generate_layout_content_slides, download_space_images
import shutil
import requests
from typing import List
from PIL import Image
from style import insert_style_slide 
from image import  SPACE_INFO

MONDAY_API_KEY = os.getenv("MONDAY_API_KEY", "YOUR_KEY_HERE")
MONDAY_API_URL = "https://api.monday.com/v2"


# -----------------------------------------------------------
#  HOTEL — SELECTED SPACES (multi_select8n5i73q8)
# -----------------------------------------------------------

def get_hotel_selected_spaces(event):
    """Extract selected hotel spaces from multi-select column."""
    col_vals = event.get("columnValues", {})
    selected_spaces = []

    if "multi_select8n5i73q8" in col_vals:
        ms = col_vals["multi_select8n5i73q8"]
        if isinstance(ms, dict) and "chosenValues" in ms:
            selected_spaces = [
                v.get("name") for v in ms["chosenValues"] if v.get("name")
            ]

    print(f"🏨 Selected hotel spaces: {selected_spaces}")
    return selected_spaces


# -----------------------------------------------------------
#  NORMALIZE SPACE NAME
# -----------------------------------------------------------

def normalize_space_name(name: str) -> str:
    return re.sub(r"\s*/\s*", "/", name.strip())


# -----------------------------------------------------------
#  FIXED: Collect ALL HOTEL pictures correctly
# -----------------------------------------------------------

def collect_all_hotel_pictures(event: dict, item_id: int) -> List[str]:
    """
    SAME LOGIC as residential:
    - use download_space_images()
    - flatten pictures
    """
    space_data = download_space_images(event, item_id)
    all_pictures = []

    for space, files in space_data.items():
        all_pictures.extend(files.get("pictures", []))

    if not all_pictures:
        print("⚠️ No HOTEL pictures found")
        return []

    print(f"🏨 Collected {len(all_pictures)} hotel pictures")
    return all_pictures


def px(value):
    return int(value * 9525)


def insert_logo_hotel_ppt(prs, event: dict):
    """
    Extract logo from files_1 (assetId), fetch public_url via Monday GraphQL,
    download the file, resize it, and insert into slide 5.
    """

    col = event.get("columnValues", {})

    # ----------------------------------------------------
    # 1️⃣ Extract assetId from files_1.files
    # ----------------------------------------------------
    files = col.get("files_1", {}).get("files", [])

    if not files:
        print("❌ No logo uploaded in files_1")
        return prs

    asset_id = files[0].get("assetId")
    print("📥 Logo asset ID:", asset_id)

    # ----------------------------------------------------
    # 2️⃣ Fetch public_url using Monday GraphQL
    # ----------------------------------------------------
    query = {
        "query": f"""
            query {{
              assets (ids: [{asset_id}]) {{
                public_url
              }}
            }}
        """
    }

    headers = {
        "Authorization": MONDAY_API_KEY,
        "Content-Type": "application/json"
    }

    try:
        response = requests.post("https://api.monday.com/v2/", json=query, headers=headers)
        response.raise_for_status()
        data = response.json()
        public_url = data["data"]["assets"][0]["public_url"]

        if not public_url:
            print("❌ public_url is missing")
            return prs

        print("🌐 public_url:", public_url)

    except Exception as e:
        print("❌ Failed retrieving public_url:", e)
        return prs

    # ----------------------------------------------------
    # 3️⃣ Download logo file
    # ----------------------------------------------------
    temp_logo = "temp_logo.png"
    resized_logo = "resized_logo.png"

    try:
        img_data = requests.get(public_url)
        img_data.raise_for_status()

        with open(temp_logo, "wb") as f:
            f.write(img_data.content)

        print("📥 Logo downloaded")

    except Exception as e:
        print("❌ Failed downloading logo:", e)
        return prs

    # ----------------------------------------------------
    # 4️⃣ Resize with correct aspect ratio
    # ----------------------------------------------------
    try:
        img = Image.open(temp_logo)
        w, h = img.size

        MAX_W = 515.9
        MAX_H = 216

        scale = min(MAX_W / w, MAX_H / h)
        new_w, new_h = int(w * scale), int(h * scale)

        img = img.resize((new_w, new_h), Image.LANCZOS)
        img.save(resized_logo)

        print(f"🖼 Logo resized to {new_w} × {new_h}px")

    except Exception as e:
        print("❌ Failed resizing logo:", e)
        return prs

    # ----------------------------------------------------
    # 5️⃣ Insert into Slide 5 at exact coordinates
    # ----------------------------------------------------
    try:
        slide = prs.slides[4]

        px_to_in = lambda px: px / 96

        X_PX = 1373.9   # ★ from your screenshot
        Y_PX = 773.8    # ★ from your screenshot

        slide.shapes.add_picture(
            resized_logo,
            Inches(px_to_in(X_PX)),
            Inches(px_to_in(Y_PX)),
            width=Inches(px_to_in(new_w)),
            height=Inches(px_to_in(new_h))
        )

        print("🎉 Logo inserted on slide 5 at correct position!")

    except Exception as e:
        print("❌ Failed inserting logo:", e)

    # ----------------------------------------------------
    # 6️⃣ Cleanup
    # ----------------------------------------------------
    for f in (temp_logo, resized_logo):
        if os.path.exists(f):
            try:
                os.remove(f)
            except:
                pass

    return prs








def insert_images_hotel_ppt(prs, image_urls: List[str]):
    """
    Insert 7 images dynamically into an existing Presentation object (prs).
    SAME LOGIC as your original version.
    FIXED to support BOTH: URLs + local file paths.
    """

    # Detect if path is URL
    def is_url(path: str) -> bool:
        return path.startswith("http://") or path.startswith("https://")

    if not image_urls:
        print("⚠️ No images provided to insert.")
        return prs

    total_needed = 7
    while len(image_urls) < total_needed:
        image_urls.extend(image_urls[: total_needed - len(image_urls)])

    temp_files = [f"temp_{i+1}.jpg" for i in range(total_needed)]
    cropped_files = [f"crop_{i+1}.jpg" for i in range(total_needed)]

    TARGET_SIZES = {
        1: (713.5, 1080),
        2: (1117, 540),
        3: (540.3, 540),
        33: (540.3, 540),
        4: (927.8, 1080),
        5: (927.8, 1080),
        6: (927.8, 1080)
    }

    def crop_to_fill(infile, outfile, target_w, target_h):
        img = Image.open(infile)
        w, h = img.size
        target_ratio = target_w / target_h
        img_ratio = w / h

        if img_ratio > target_ratio:
            new_w = int(h * target_ratio)
            left = (w - new_w) // 2
            img = img.crop((left, 0, left + new_w, h))
        else:
            new_h = int(w / target_ratio)
            top = (h - new_h) // 2
            img = img.crop((0, top, w, top + new_h))

        img = img.resize((int(target_w), int(target_h)), Image.LANCZOS)
        img = img.convert("RGB")
        img.save(outfile, "JPEG")

    # ------------------------------------------------------------------
    # 🔥 DOWNLOAD OR COPY LOCAL FILES + CROP
    # ------------------------------------------------------------------
    for i, src in enumerate(image_urls[:total_needed]):
        try:
            temp_path = temp_files[i]

            if is_url(src):
                print(f"🌐 Downloading URL: {src}")
                response = requests.get(src, timeout=10)
                response.raise_for_status()
                with open(temp_path, "wb") as f:
                    f.write(response.content)
            else:
                print(f"📁 Using local image: {src}")
                if not os.path.exists(src):
                    print(f"❌ Local file not found: {src}")
                    continue
                shutil.copy(src, temp_path)

            # pick target box
            if i == 0:
                tw, th = TARGET_SIZES[1]
            elif i == 1:
                tw, th = TARGET_SIZES[2]
            elif i == 2:
                tw, th = TARGET_SIZES[3]
            elif i == 3:
                tw, th = TARGET_SIZES[33]
            elif i == 4:
                tw, th = TARGET_SIZES[4]
            elif i == 5:
                tw, th = TARGET_SIZES[5]
            elif i == 6:
                tw, th = TARGET_SIZES[6]

            crop_to_fill(temp_path, cropped_files[i], tw, th)
            print(f"✅ Cropped {cropped_files[i]}")

        except Exception as e:
            print(f"⚠️ Failed processing image {src}: {e}")

    # ------------------------------------------------------------------
    # INSERT BEHIND CONTENT
    # ------------------------------------------------------------------
    def insert_picture_behind(slide, image_path, left, top, width, height):
        pic = slide.shapes.add_picture(image_path, left, top, width, height)
        spTree = slide.shapes._spTree
        spTree.remove(pic._element)
        spTree.insert(2, pic._element)

    print("🖼 Inserting images into Hotel cover slides...")

    insert_picture_behind(prs.slides[0], cropped_files[0], px(
        164), px(0), px(713.5), px(1080))
    insert_picture_behind(prs.slides[1], cropped_files[1], px(
        803), px(0), px(1117), px(540))
    insert_picture_behind(prs.slides[2], cropped_files[2], px(
        596.9), px(0), px(540.3), px(540))
    insert_picture_behind(prs.slides[2], cropped_files[3], px(
        1137.2), px(0), px(540.3), px(540))
    insert_picture_behind(prs.slides[3], cropped_files[4], px(
        165.4), px(0), px(927.8), px(1080))
    insert_picture_behind(prs.slides[4], cropped_files[5], px(
        165.4), px(0), px(927.8), px(1080))
    insert_picture_behind(prs.slides[5], cropped_files[6], px(
        165.4), px(0), px(927.8), px(1080))

    # CLEANUP
    for file_path in temp_files + cropped_files:
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
            except:
                pass

    print("🧹 Temporary images removed successfully.")
    return prs


def safe_folder_name(name: str):
    name = name.strip().replace("/", "_").replace("\\", "_")
    for ch in '<>:"/\\|?*':
        name = name.replace(ch, "_")
    return name


# -----------------------------------------------------------
#  BUILD HOTEL SPACE PPT — FIXED + SAFE
# -----------------------------------------------------------

def build_hotel_space_ppt(template_ppt_path, space_data, prs, selected_style, event):
    """
    Inserts slides for each hotel space into prs.
    """
    insert_index = 6  # hotel cover uses 6 slides

    for space, files in space_data.items():
        print(f"\n📂 Building slides for hotel space: {space}")

        pictures = files.get("pictures", [])
        floorplans = files.get("floor_plans", [])
        elevations = files.get("elevations", [])
        inspiration = files.get("inspiration", [])

        # ---- Inspiration ----
        if inspiration:
            print("✨ Creating Inspiration slides...")
            insert_index = generate_inspiration_slides(
                "templates/inspiration_slides_template.pptx",
                inspiration,
                prs,
                insert_index
            )
          # ⭐ NEW: Insert STYLE SLIDE AFTER INSPIRATION

        # ---- Floorplans + Pictures ----
        if pictures and floorplans:
            insert_index = generate_layout_content_slides(
                template_ppt_path,
                floorplans,
                pictures,
                prs,
                insert_index,
                event
            )

        # ---- Elevations + Pictures ----
        if pictures and elevations:
            insert_index = generate_layout_content_slides(
                template_ppt_path,
                elevations,
                pictures,
                prs,
                insert_index,
                event
            )

        # ---- Pictures only ----
        if pictures and not (floorplans or elevations):
            insert_index = generate_layout_content_slides(
                template_ppt_path,
                [pictures[0]],
                pictures,
                prs,
                insert_index,
                event
            )
    if selected_style:
        insert_index = insert_style_slide(
            prs,
            selected_style,  # handles list internally
            insert_index
        )

    return prs, insert_index

def compute_combined_area(col_values):
    """
    Reads selected spaces from dropdown0,
    fetches their area fields, and returns TOTAL area:
        e.g., '80 m2'
    """
    selected = col_values.get("dropdown0", {}).get("chosenValues", [])
    selected_spaces = [item["name"] for item in selected]

    if not selected_spaces:
        return ""

    total_area = 0

    for space in selected_spaces:
        if space not in SPACE_INFO:
            continue

        area_col_id = SPACE_INFO[space]["total_area"]
        area_col = col_values.get(area_col_id, {})

        raw_val = area_col.get("value") or area_col.get("text") or ""

        # Ensure numeric
        if raw_val.strip().isdigit():
            total_area += int(raw_val)

    if total_area == 0:
        return ""

    return f"{total_area} m2"

def clean_styles(raw_styles):
    """Convert style field into a clean, comma-separated string."""
    if not raw_styles:
        return ""

    # If Typeform returns as list
    if isinstance(raw_styles, list):
        return ", ".join(s.strip() for s in raw_styles if s.strip())

    # If it is a single string, normalize separators
    text = str(raw_styles).strip()

    # Replace multiple separators with commas
    for sep in [";", "|", "/", "  "]:
        text = text.replace(sep, ",")

    # Split on commas and clean up
    parts = [p.strip() for p in text.split(",") if p.strip()]

    return ", ".join(parts)


# -----------------------------------------------------------
#  FORM MAPPING (unchanged)
# -----------------------------------------------------------

def map_hotel_form_data(event: dict):
    col = event.get("columnValues", {})
    styles = []

    # --- STYLE HANDLING ---
    style_columns = ["dropdown", "dropdown0",
                     "dropdown1", "dropdown2", "style_dropdown"]
    for col_name in style_columns:
        if col_name in col and col[col_name].get("chosenValues"):
            chosen_values = col[col_name]["chosenValues"]
            styles = [v.get("name", "")
                      for v in chosen_values if v.get("name")]
            break

    styles_text = ",".join(styles) if styles else ""
    styles = clean_styles(styles_text)

    combined_area = compute_combined_area(col)

    return {
        "Q. Nature of the project": col.get("status1", {}).get("label", {}).get("text"),
        "Q. Project Name": col.get("text_10", {}).get("value"),
        "city": col.get("text8", {}).get("value"),
        "9. Project address": col.get("short_text9xa0p91x", {}).get("value"),
        "Type of project": col.get("dropdown76", {}).get("chosenValues", [{}])[0].get("name"),
        # Additional calculated area
        "Q. Area": combined_area,

        # ✅ FIXED FIELD — uses multi_select8n5i73q8 now
        "What space(s) are you looking to design":
            ", ".join([v.get("name", "") for v in col.get(
                "multi_select8n5i73q8", {}).get("chosenValues", [])]),

        "7. Website": col.get("text96", {}).get("value"),
        "10. Planned or current Category": col.get("single_selectx3pcawh", {}).get("label", {}).get("text"),
        "11. Targeted Clientele":
            ", ".join([v.get("name", "") for v in col.get(
                "multi_select64xarcdd", {}).get("chosenValues", [])]),

        "13. Number of guest rooms planned or existing": col.get("short_textq3gq73ug", {}).get("value"),
        "3. Your Project / Hotel Name": col.get("text_10", {}).get("value"),

        "18. What are the main objectives of the project":
            ", ".join([v.get("name", "") for v in col.get(
                "multi_selectbw4ovak9", {}).get("chosenValues", [])]),

        "19. Estimated budget for FF&E + works":
            col.get("single_selectye6dyr5", {}).get("label", {}).get("text"),

        "20. Planned opening date or deadline": col.get("date0r1120f0", {}).get("date"),

        "21. Which constraints":
            ", ".join([v.get("name", "") for v in col.get(
                "multi_selectgin5h4yr", {}).get("chosenValues", [])]),

        "31. Other comments or details": col.get("long_text4n15rau0", {}).get("text"),

        # --- DESIGN SECTION ---
        " 22. Design Guideine":
            col.get("single_select2ttpyke", {}).get("label", {}).get("text"),


        "24. Do you have a defined color palette":
            col.get("single_selectdgir7ru", {}).get("label", {}).get("text"),

        "25. Please describe the style / atmosphere":
            col.get("long_textcl38cdjs", {}).get("text"),

        "29. How important is sustainability in your project":
            col.get("single_selectiz3tpad", {}).get("label", {}).get("text"),

        "28. Would you like the design to reflect local culture":
            col.get("single_selectexftwo4", {}).get("label", {}).get("text"),

        # ✔️ Correct for 27. Style
        "Which style's do you like1": styles,

        
    }





